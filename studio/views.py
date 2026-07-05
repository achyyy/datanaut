import os
import sqlite3
import pandas as pd
from django.shortcuts import render, redirect
from django.http import JsonResponse, HttpResponse
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
from django.urls import reverse
from groq import Groq
from dotenv import load_dotenv
from django.contrib import messages
import time
from datetime import datetime, timezone, timedelta
from xhtml2pdf import pisa
import json
import pdfkit, io
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
import io
from django.conf import settings
from supabase import create_client
import base64
from io import BytesIO
import requests


# Load environment variables
load_dotenv()
print(f"DEBUG MODULE INIT: Loading GROQ_API_KEY...")
groq_key = os.getenv("GROQ_API_KEY")
print(f"DEBUG MODULE INIT: GROQ_API_KEY loaded: {groq_key[:20] if groq_key else 'NOT FOUND'}...")
print(f"DEBUG MODULE INIT: GROQ_API_KEY length: {len(groq_key) if groq_key else 0}")
client = Groq(api_key=groq_key)
print(f"DEBUG MODULE INIT: Groq client initialized successfully")

# Supabase client
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")
SUPABASE_BUCKET = os.getenv("SUPABASE_BUCKET")
supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)

TEMP_FILES = {}

def get_active_dataset_path(request):
    """
    Helper function to get the active dataset path from the multi-dataset session.
    This should be used instead of request.session.get("dataset_path") everywhere.
    
    Returns:
        tuple: (supabase_path, display_name) or (None, None) if no dataset found
    """
    datasets_info = request.session.get('datasets', [])
    active_index = request.session.get('active_dataset_index', 0)
    
    if not datasets_info:
        # Fallback to old single-dataset system for backward compatibility
        old_path = request.session.get('current_dataset_path')
        old_name = request.session.get('dataset_display_name', 'Dataset')
        return old_path, old_name
    
    if 0 <= active_index < len(datasets_info):
        active_dataset = datasets_info[active_index]
        return active_dataset['supabase_path'], active_dataset['display_name']
    
    return None, None

def cleanup_session_data_from_supabase(session):
    """Delete all session-related data from Supabase storage"""
    try:
        # Get session-specific files to delete
        session_files = []
        
        # Check if there's a current dataset to delete
        if 'current_dataset_path' in session:
            dataset_path = session['current_dataset_path']
            # Extract filename from path for Supabase
            if '/' in dataset_path:
                filename = dataset_path.split('/')[-1]
                session_files.append(f"sessions/{filename}")
        
        # Check for multi-dataset array
        if 'datasets' in session:
            for dataset_info in session['datasets']:
                if isinstance(dataset_info, dict) and 'supabase_path' in dataset_info:
                    session_files.append(dataset_info['supabase_path'])
        
        # Also check for any session-specific uploaded files (legacy)
        if 'uploaded_files' in session:
            for file_info in session['uploaded_files']:
                if isinstance(file_info, dict) and 'supabase_path' in file_info:
                    session_files.append(file_info['supabase_path'])
                elif isinstance(file_info, str):
                    session_files.append(f"sessions/{file_info}")
        
        # Delete files from Supabase
        if session_files:
            for file_path in session_files:
                try:
                    supabase.storage.from_(SUPABASE_BUCKET).remove([file_path])
                    print(f"✅ Deleted from Supabase: {file_path}")
                except Exception as file_error:
                    print(f"⚠️ Failed to delete {file_path}: {file_error}")
        
        print(f"🧹 Cleaned up {len(session_files)} files from Supabase")
        
    except Exception as e:
        print(f"❌ Error cleaning up Supabase data: {e}")

def complete_session_cleanup(session):
    """Complete cleanup of session data including Supabase storage"""
    try:
        # First, clean up Supabase data
        cleanup_session_data_from_supabase(session)
        
        # Then clear all session data
        session_keys_to_clear = [
            'chat_history', 'visualizations', 'tables', 'report_blocks', 
            'query_results', 'current_dataset_path', 'uploaded_files',
            'dataset_info', 'analysis_results', 'generated_charts',
            'ai_responses', 'user_queries', 'session_start_time',
            'datasets', 'active_dataset_index'  # Multi-dataset keys
        ]
        
        cleared_count = 0
        for key in session_keys_to_clear:
            if key in session:
                del session[key]
                cleared_count += 1
        
        # Force session save
        session.modified = True
        session.save()
        
        print(f"🧹 Completed session cleanup - cleared {cleared_count} session variables")
        return True
        
    except Exception as e:
        print(f"❌ Error in complete session cleanup: {e}")
        return False

def cleanup_old_files():
    try:
        files = supabase.storage.from_(SUPABASE_BUCKET).list("sessions/")
        for f in files:
            created_at = datetime.fromisoformat(f["created_at"].replace("Z", "+00:00"))
            now = datetime.now(timezone.utc)
            if now - created_at > timedelta(hours=1):
                supabase.storage.from_(SUPABASE_BUCKET).remove([f["name"]])
                print(f"✅ Deleted old file from Supabase: {f['name']}")

    except Exception as e:
        print(f"⚠️ Cleanup failed: {e}")


def clear_session_data(request):
    """Clear all previous session data when uploading new dataset"""
    keys_to_clear = [
        'chat_history',
        'visualizations', 
        'tables',
        'report_blocks',
        'query_results'
    ]
    
    for key in keys_to_clear:
        if key in request.session:
            del request.session[key]
    
    print("✅ Cleared previous session data for fresh start")

def index(request):
    return render(request, "index.html", {"datasets": []})


def search_online_datasets(request):
    """AI-powered dataset search engine"""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'})
    
    try:
        data = json.loads(request.body)
        search_query = data.get('query', '').strip()
        
        if not search_query:
            return JsonResponse({'status': 'error', 'message': 'Please enter a search query'})
        
        print(f"🔍 Searching for datasets: {search_query}")
        
        # Use AI to find relevant datasets
        prompt = f"""You are a dataset search engine. Find relevant CSV datasets for this query: "{search_query}"

Generate 5-8 dataset suggestions using ONLY these verified sources:

VERIFIED SOURCES:
1. https://raw.githubusercontent.com/mwaskom/seaborn-data/master/[filename].csv
   - Available: iris.csv, titanic.csv, tips.csv, planets.csv, diamonds.csv, flights.csv, penguins.csv

2. https://raw.githubusercontent.com/datasciencedojo/datasets/master/[filename].csv
   - Available: titanic.csv, iris.csv, Admission_Predict.csv, Mall_Customers.csv

3. https://people.sc.fsu.edu/~jburkardt/data/csv/[filename].csv
   - Available: airtravel.csv, biostats.csv, cities.csv, deniro.csv, hw_200.csv, zillow.csv

4. https://raw.githubusercontent.com/plotly/datasets/master/[filename].csv
   - Available: iris.csv, wind_data.csv, 2011_us_ag_exports.csv, volcano.csv

INSTRUCTIONS:
- Use ONLY URLs from the verified sources above
- Match the search query to the most relevant available datasets
- Provide accurate descriptions based on well-known datasets
- If no exact match, suggest similar datasets from the list

Return ONLY a JSON array (no markdown, no explanation):
[
  {{
    "name": "Dataset Name",
    "description": "What the dataset contains",
    "category": "Category",
    "rows": "estimated number",
    "columns": estimated_number,
    "url": "https://raw.githubusercontent.com/..."
  }}
]"""

        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1500
        )
        
        response_text = response.choices[0].message.content.strip()
        
        # Extract JSON from response (handle markdown code blocks)
        if '```json' in response_text:
            response_text = response_text.split('```json')[1].split('```')[0].strip()
        elif '```' in response_text:
            response_text = response_text.split('```')[1].split('```')[0].strip()
        
        # Parse the JSON
        datasets = json.loads(response_text)
        
        # Add IDs
        for i, ds in enumerate(datasets):
            ds['id'] = f"ai_dataset_{i}_{int(time.time())}"
            if 'size' not in ds:
                ds['size'] = 'Unknown'
        
        print(f"✅ Found {len(datasets)} datasets for: {search_query}")
        
        return JsonResponse({
            'status': 'success',
            'datasets': datasets,
            'query': search_query
        })
        
    except json.JSONDecodeError as e:
        print(f"❌ JSON parsing error: {e}")
        print(f"Response was: {response_text}")
        return JsonResponse({'status': 'error', 'message': 'Failed to parse AI response'})
    except Exception as e:
        print(f"❌ Dataset search error: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'status': 'error', 'message': str(e)})


def fetch_online_dataset(request):
    """Fetch an online dataset and save it to session"""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'})
    
    try:
        import requests
        data = json.loads(request.body)
        dataset_id = data.get('dataset_id')
        dataset_url = data.get('url')
        dataset_name = data.get('name', 'Online Dataset')
        
        if not dataset_url:
            return JsonResponse({'status': 'error', 'message': 'Dataset URL is required'})
        
        print(f"🌐 Fetching online dataset: {dataset_name}")
        print(f"📍 URL: {dataset_url}")
        
        # Fetch the dataset from URL with headers to avoid blocks
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            'Accept': 'text/csv,text/plain,application/csv,*/*'
        }
        
        response = requests.get(dataset_url, timeout=30, headers=headers, allow_redirects=True)
        response.raise_for_status()
        
        # Verify it's CSV data
        content = response.content
        if len(content) < 10:
            return JsonResponse({'status': 'error', 'message': 'Dataset appears to be empty'})
        
        # Try to parse as CSV to validate
        try:
            test_df = pd.read_csv(io.BytesIO(content))
            if test_df.empty:
                return JsonResponse({'status': 'error', 'message': 'Dataset is empty'})
            print(f"✅ Validated CSV: {test_df.shape[0]} rows, {test_df.shape[1]} columns")
        except Exception as e:
            print(f"⚠️ CSV validation warning: {e}")
            # Try to continue anyway
        
        # Save to Supabase (append to existing datasets, don't clear)
        uploaded_datasets = request.session.get('datasets', [])
        
        session_key = request.session.session_key or str(int(time.time()))
        timestamp = int(time.time() * 1000)  # Use milliseconds for consistency
        
        # Clean filename
        safe_name = dataset_name.replace(' ', '_').replace('/', '_')[:50]
        stored_name = f"{session_key}_{timestamp}_{safe_name}.csv"
        path_in_bucket = f"sessions/{stored_name}"
        
        # Upload to Supabase
        supabase.storage.from_(SUPABASE_BUCKET).upload(
            path_in_bucket,
            content,
            {"cacheControl": "3600"}
        )
        
        print(f"✅ Uploaded online dataset to Supabase: {path_in_bucket}")
        
        # Append to existing datasets (not replace)
        uploaded_datasets.append({
            'supabase_path': path_in_bucket,
            'stored_name': stored_name,
            'display_name': dataset_name if dataset_name.endswith('.csv') else f"{dataset_name}.csv",
            'upload_time': timestamp,
            'size': len(content),
            'is_online': True,
            'source_url': dataset_url
        })
        
        request.session['datasets'] = uploaded_datasets
        request.session['active_dataset_index'] = len(uploaded_datasets) - 1
        request.session['current_dataset_path'] = path_in_bucket
        request.session['uploaded_files'] = request.session['datasets']
        request.session.modified = True
        
        messages.success(request, f"Successfully loaded online dataset: {dataset_name}")
        
        return JsonResponse({
            'status': 'success',
            'message': f'Loaded {dataset_name}',
            'redirect_url': reverse("multi_dataset_preview")
        })
        
    except requests.exceptions.Timeout:
        print(f"❌ Timeout fetching dataset")
        return JsonResponse({'status': 'error', 'message': 'Request timeout - dataset may be too large or server slow'})
    except requests.exceptions.RequestException as e:
        print(f"❌ Failed to fetch online dataset: {e}")
        return JsonResponse({'status': 'error', 'message': f'Failed to download: {str(e)}'})
    except Exception as e:
        print(f"❌ Error fetching online dataset: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'status': 'error', 'message': f'Error: {str(e)}'})


def dataset_list(request):
    return render(request, "datasets.html", {"datasets": []})


def upload_dataset(request):
    """Legacy single file upload - redirects to multi-upload"""
    if request.method == "POST" and request.FILES.get("file"):
        file = request.FILES["file"]
        request.FILES.setlist('files', [file])
        return upload_multiple_datasets(request)
    return redirect("index")


def upload_multiple_datasets(request):
    """Handle multiple dataset uploads"""
    if request.method == "POST":
        files = request.FILES.getlist("files")
        
        if not files:
            return JsonResponse({'status': 'error', 'message': 'No files provided'})
        
        # Ensure session exists
        if not request.session.session_key:
            request.session.create()
        
        session_key = request.session.session_key
        print(f"📝 Session key: {session_key}")
        
        # Get existing datasets instead of clearing them (append mode)
        uploaded_datasets = request.session.get('datasets', [])
        request.session['uploaded_files'] = request.session.get('uploaded_files', [])
        
        for file in files:
            try:
                timestamp = int(time.time() * 1000)  # Use milliseconds for uniqueness
                stored_name = f"{session_key}_{timestamp}_{file.name}"
                path_in_bucket = f"sessions/{stored_name}"
                display_name = file.name
                
                print(f"DEBUG: Uploading file: {stored_name}")
                
                # Read file data
                data = file.read()
                
                # Upload to Supabase (with upsert to handle duplicates)
                try:
                    supabase.storage.from_(SUPABASE_BUCKET).upload(
                        path_in_bucket,
                        data,
                        {"cacheControl": "3600", "upsert": "true"}
                    )
                    print(f"✅ Uploaded to Supabase: {path_in_bucket}")
                except Exception as upload_err:
                    # If upload fails, try to remove and re-upload
                    print(f"⚠️ Upload conflict, attempting to replace: {upload_err}")
                    try:
                        supabase.storage.from_(SUPABASE_BUCKET).remove([path_in_bucket])
                        supabase.storage.from_(SUPABASE_BUCKET).upload(
                            path_in_bucket,
                            data,
                            {"cacheControl": "3600"}
                        )
                        print(f"✅ Re-uploaded to Supabase: {path_in_bucket}")
                    except Exception as retry_err:
                        print(f"❌ Failed to upload after retry: {retry_err}")
                        raise
                
                # Store dataset info
                dataset_info = {
                    'supabase_path': path_in_bucket,
                    'stored_name': stored_name,
                    'display_name': display_name,
                    'upload_time': timestamp,
                    'size': len(data)
                }
                
                uploaded_datasets.append(dataset_info)
                print(f"   ✅ Added to uploaded_datasets: {display_name}")
                
                # Small delay to ensure unique timestamps
                time.sleep(0.01)
                
            except Exception as e:
                print(f"❌ Failed to upload {file.name}: {e}")
                # Continue with other files
                continue
        
        if not uploaded_datasets:
            return JsonResponse({'status': 'error', 'message': 'All uploads failed'})
        
        # Store all datasets in session
        request.session['datasets'] = uploaded_datasets
        request.session['uploaded_files'] = uploaded_datasets  # Keep in sync
        request.session['current_dataset_path'] = uploaded_datasets[-1]['supabase_path']
        request.session['active_dataset_index'] = len(uploaded_datasets) - 1  # Set active to the last uploaded file
        request.session.modified = True
        
        print(f"\n📦 Session state after upload:")
        print(f"   Total datasets: {len(uploaded_datasets)}")
        print(f"   Active index: 0")
        print(f"   Datasets: {[d['display_name'] for d in uploaded_datasets]}")
        
        messages.success(request, f"Successfully uploaded {len(uploaded_datasets)} dataset(s)")
        
        # Return JSON response for AJAX or redirect for form submission
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest' or request.content_type == 'multipart/form-data':
            return JsonResponse({
                'status': 'success',
                'message': f'Uploaded {len(uploaded_datasets)} files',
                'redirect_url': reverse("multi_dataset_preview"),
                'datasets': uploaded_datasets
            })
        else:
            return redirect("multi_dataset_preview")
    
    return redirect("index")


def multi_dataset_preview(request):
    """Preview and manage multiple datasets with tabs"""
    datasets_info = request.session.get('datasets', [])
    
    if not datasets_info:
        messages.error(request, "No datasets in session.")
        return redirect("index")
    
    # Clamp active index to valid range to avoid IndexError
    active_index = request.session.get('active_dataset_index', 0)
    if not isinstance(active_index, int):
        active_index = 0
    if datasets_info:
        max_idx = len(datasets_info) - 1
        if active_index < 0 or active_index > max_idx:
            active_index = 0
            request.session['active_dataset_index'] = active_index
            request.session.modified = True
    
    # Load all datasets from Supabase
    loaded_datasets = []
    for idx, dataset_info in enumerate(datasets_info):
        try:
            print(f"🔍 DEBUG: Loading dataset {idx}: {dataset_info['display_name']}")
            print(f"   Path: {dataset_info['supabase_path']}")
            
            res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_info['supabase_path'])
            
            print(f"   Response type: {type(res)}, size: {len(res) if res else 0}")
            
            df = pd.read_csv(io.BytesIO(res))
            
            # Generate profile for each dataset with JSON-safe data
            profile = {
                'index': idx,
                'display_name': dataset_info['display_name'],
                'stored_name': dataset_info['stored_name'],
                'shape': list(df.shape),  # Convert tuple to list
                'columns': df.columns.tolist(),
                'preview_rows': sanitize_for_json(df.head(5).values.tolist()),
                'dtypes': {col: str(dtype) for col, dtype in df.dtypes.items()},
                'is_active': idx == active_index
            }
            
            loaded_datasets.append(profile)
            print(f"   ✅ Successfully loaded: {df.shape}")
            
        except Exception as e:
            print(f"❌ Failed to load dataset {dataset_info['display_name']}: {e}")
            print(f"   Dataset info: {dataset_info}")
            
            # Try to list available files to help debug
            try:
                all_files = supabase.storage.from_(SUPABASE_BUCKET).list('sessions/')
                available_files = [f['name'] for f in all_files if f.get('name')]
                print(f"   Available files in sessions/: {available_files}")
                
                # Check if a similar file exists (with different timestamp)
                expected_basename = dataset_info['display_name']
                similar_files = [f for f in available_files if expected_basename in f]
                if similar_files:
                    print(f"   💡 Similar files found: {similar_files}")
                    print(f"   💡 TIP: Your session data is outdated. Please re-upload your files or clear session.")
            except Exception as list_err:
                print(f"   Could not list files: {list_err}")
            
            import traceback
            traceback.print_exc()
            continue
    
    if not loaded_datasets:
        messages.error(request, "Failed to load datasets.")
        return redirect("index")
    
    # Generate insights for active dataset only
    try:
        active_dataset_info = datasets_info[active_index]
        res = supabase.storage.from_(SUPABASE_BUCKET).download(active_dataset_info['supabase_path'])
        active_df = pd.read_csv(io.BytesIO(res))
        
        data_profile = generate_data_profile(active_df)
        analysis = analyze_dataset_context(active_df)
        curious_questions = generate_curious_insights(active_df)
        
    except Exception as e:
        print(f"❌ Failed to analyze active dataset: {e}")
        data_profile = None
        analysis = None
        curious_questions = []
    
    # Select the loaded profile that matches the active index; fallback to first
    active_loaded = None
    if loaded_datasets:
        active_loaded = next((d for d in loaded_datasets if d.get('index') == active_index), loaded_datasets[0])

    context = {
        'datasets': loaded_datasets,
        'active_dataset': active_loaded,
        'data_profile': data_profile,
        'analysis': analysis,
        'curious_questions': curious_questions,
        'can_merge': len(loaded_datasets) > 1
    }
    
    return render(request, "multi_dataset_preview.html", context)


def switch_dataset(request, dataset_index):
    """Switch active dataset"""
    datasets_info = request.session.get('datasets', [])
    
    print(f"🔄 Switch dataset requested:")
    print(f"   Requested index: {dataset_index}")
    print(f"   Total datasets: {len(datasets_info)}")
    print(f"   Available datasets: {[d['display_name'] for d in datasets_info]}")
    
    if 0 <= dataset_index < len(datasets_info):
        request.session['active_dataset_index'] = dataset_index
        active_dataset = datasets_info[dataset_index]
        # Update backward compatibility key
        request.session['current_dataset_path'] = active_dataset['supabase_path']
        request.session.modified = True
        
        print(f"   ✅ Switched to: {active_dataset['display_name']}")
        
        return JsonResponse({
            'status': 'success', 
            'active_index': dataset_index,
            'dataset_name': active_dataset['display_name']
        })
    
    print(f"   ❌ Invalid index")
    return JsonResponse({'status': 'error', 'message': 'Invalid dataset index'})


def merge_datasets(request):
    """Merge multiple datasets based on user specifications"""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'})
    
    try:
        data = json.loads(request.body)
        merge_type = data.get('merge_type', 'concat')  # concat, join
        dataset_indices = data.get('dataset_indices', [])
        
        datasets_info = request.session.get('datasets', [])
        
        if len(dataset_indices) < 2:
            return JsonResponse({'status': 'error', 'message': 'Select at least 2 datasets to merge'})
        
        # Load selected datasets
        dfs = []
        names = []
        for idx in dataset_indices:
            if 0 <= idx < len(datasets_info):
                dataset_info = datasets_info[idx]
                res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_info['supabase_path'])
                df = pd.read_csv(io.BytesIO(res))
                dfs.append(df)
                names.append(dataset_info['display_name'])
        
        if len(dfs) < 2:
            return JsonResponse({'status': 'error', 'message': 'Failed to load datasets'})
        
        # Perform merge based on type
        if merge_type == 'concat':
            # Vertical concatenation (stack rows)
            merged_df = pd.concat(dfs, ignore_index=True)
            merge_desc = f"Concatenated {len(dfs)} datasets vertically"
            
        elif merge_type == 'join':
            # Horizontal join (need common column)
            join_column = data.get('join_column')
            join_type = data.get('join_type', 'inner')  # inner, outer, left, right
            
            if not join_column:
                return JsonResponse({'status': 'error', 'message': 'Join column required for join operation'})
            
            merged_df = dfs[0]
            for i in range(1, len(dfs)):
                merged_df = merged_df.merge(dfs[i], on=join_column, how=join_type, suffixes=('', f'_{names[i]}'))
            
            merge_desc = f"Joined {len(dfs)} datasets on column '{join_column}' using {join_type} join"
        
        else:
            return JsonResponse({'status': 'error', 'message': 'Invalid merge type'})
        
        # Save merged dataset to Supabase
        session_key = request.session.session_key or str(int(time.time()))
        timestamp = int(time.time() * 1000)  # Use milliseconds for consistency
        merged_name = f"{session_key}_{timestamp}_merged.csv"
        path_in_bucket = f"sessions/{merged_name}"
        
        # Convert to CSV
        csv_buffer = io.StringIO()
        merged_df.to_csv(csv_buffer, index=False)
        csv_data = csv_buffer.getvalue().encode('utf-8')
        
        # Upload to Supabase
        supabase.storage.from_(SUPABASE_BUCKET).upload(
            path_in_bucket,
            csv_data,
            {"cacheControl": "3600"}
        )
        
        # Add merged dataset to session
        merged_info = {
            'supabase_path': path_in_bucket,
            'stored_name': merged_name,
            'display_name': f"Merged_{'+'.join(names[:2])}",
            'upload_time': timestamp,
            'size': len(csv_data),
            'is_merged': True
        }
        
        datasets_info.append(merged_info)
        request.session['datasets'] = datasets_info
        request.session['active_dataset_index'] = len(datasets_info) - 1
        request.session['uploaded_files'].append(merged_info)
        request.session.modified = True
        
        return JsonResponse({
            'status': 'success',
            'message': merge_desc,
            'merged_dataset': {
                'name': merged_info['display_name'],
                'shape': merged_df.shape,
                'columns': merged_df.columns.tolist()
            }
        })
        
    except Exception as e:
        print(f"❌ Merge failed: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'status': 'error', 'message': str(e)})


def run_multi_dataset_query(request):
    """Run queries across multiple datasets"""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'})
    
    try:
        data = json.loads(request.body)
        query = data.get('query', '')
        dataset_index = data.get('dataset_index')
        
        datasets_info = request.session.get('datasets', [])
        
        if dataset_index is None:
            dataset_index = request.session.get('active_dataset_index', 0)
        
        if not (0 <= dataset_index < len(datasets_info)):
            return JsonResponse({'status': 'error', 'message': 'Invalid dataset'})
        
        # Load the selected dataset
        dataset_info = datasets_info[dataset_index]
        print(f"🔍 QUERY DEBUG: Loading dataset {dataset_index}: {dataset_info['display_name']}")
        print(f"   Path: {dataset_info['supabase_path']}")
        
        res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_info['supabase_path'])
        print(f"   Downloaded {len(res)} bytes")
        
        df = pd.read_csv(io.BytesIO(res))
        print(f"   Loaded dataframe: {df.shape}")
        
        # Get column info
        columns = [{'name': col, 'type': str(df[col].dtype)} for col in df.columns]
        
        # Generate SQL using AI
        prompt = make_sql_prompt(dataset_info['display_name'], columns, query)
        
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            max_tokens=500
        )
        
        # Extract SQL from AI response
        raw_response = response.choices[0].message.content.strip()
        parsed = extract_top_level_json(raw_response)
        sql_query = parsed.get("sql", "").strip()
        
        if not sql_query:
            raise ValueError("AI did not return valid SQL query")
        
        # Execute query
        conn = sqlite3.connect(":memory:")
        df.to_sql(dataset_info['display_name'], conn, index=False, if_exists="replace")
        
        result_df = pd.read_sql_query(sql_query, conn)
        conn.close()
        
        # Format response with sanitized data
        result_data = sanitize_for_json(result_df.to_dict('records'))

        # ---- Persist chat + query to session for Enhanced Report Builder ----
        try:
            # Save the user query as a chat message
            save_chat_message(request, "User", query)

            # Save AI response summary (SQL + rows)
            ai_response = f"Generated SQL Query:\n```sql\n{sql_query}\n```\n\nFound {len(result_data)} results from {dataset_info['display_name']}."
            save_chat_message(request, "AI Assistant", ai_response)

            # Save full query details for table rendering in reports
            save_query_result(
                request,
                query,
                sql_query,
                result_data,
                suggest_visualization_type(result_df, sql_query)
            )
        except Exception as persist_err:
            print(f"⚠️ Failed to persist chat/query to session: {persist_err}")
        
        return JsonResponse({
            'status': 'success',
            'query': query,
            'sql': sql_query,
            'results': result_data,
            'columns': result_df.columns.tolist(),
            'row_count': len(result_df),
            'dataset_name': dataset_info['display_name']
        })
        
    except Exception as e:
        print(f"❌ Query failed: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'status': 'error', 'message': str(e)})



def generate_data_profile(df):
    """Generate comprehensive data profile for health checks"""
    profile = {
        'shape': df.shape,
        'columns': [],
        'data_health': {
            'duplicate_rows': df.duplicated().sum(),
            'total_missing': df.isnull().sum().sum(),
            'missing_percentage': (df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100,
            'constant_columns': [],
            'high_cardinality_columns': []
        },
        'column_stats': {}
    }
    
    for col in df.columns:
        col_info = {
            'name': str(col),
            'dtype': str(df[col].dtype),
            'missing_count': int(df[col].isnull().sum()),
            'missing_percentage': float((df[col].isnull().sum() / len(df)) * 100),
            'unique_count': int(df[col].nunique()),
            'cardinality_ratio': float(df[col].nunique() / len(df)),
            'sample_values': [str(x) for x in df[col].dropna().head(3).tolist()] if len(df[col].dropna()) > 0 else []
        }
        
        # Check if column is constant
        if col_info['unique_count'] <= 1:
            profile['data_health']['constant_columns'].append(col)
        
        # Check for high cardinality (might be ID columns)
        if col_info['cardinality_ratio'] > 0.9 and col_info['unique_count'] > 10:
            profile['data_health']['high_cardinality_columns'].append(col)
        
        # Type-specific stats
        if df[col].dtype in ['int64', 'float64']:
            col_info.update({
                'min': float(df[col].min()) if not df[col].isnull().all() else None,
                'max': float(df[col].max()) if not df[col].isnull().all() else None,
                'mean': float(df[col].mean()) if not df[col].isnull().all() else None,
                'median': float(df[col].median()) if not df[col].isnull().all() else None,
                'std': float(df[col].std()) if not df[col].isnull().all() else None,
                'has_outliers': bool(detect_outliers(df[col])) if not df[col].isnull().all() else False
            })
        elif df[col].dtype == 'object':
            if not df[col].isnull().all():
                value_counts = df[col].value_counts()
                col_info.update({
                    'top_values': value_counts.head(5).to_dict(),
                    'is_categorical': col_info['unique_count'] < len(df) * 0.5,
                    'avg_length': df[col].astype(str).str.len().mean()
                })
        
        profile['columns'].append(col_info)
        profile['column_stats'][col] = col_info
    
    # Convert all numpy/pandas types to JSON serializable types
    profile = make_json_serializable(profile)
    
    return profile

def detect_outliers(series):
    """Simple outlier detection using IQR method"""
    try:
        Q1 = series.quantile(0.25)
        Q3 = series.quantile(0.75)
        IQR = Q3 - Q1
        lower_bound = Q1 - 1.5 * IQR
        upper_bound = Q3 + 1.5 * IQR
        outliers = series[(series < lower_bound) | (series > upper_bound)]
        return len(outliers) > 0
    except:
        return False

def make_json_serializable(obj):
    """Convert numpy/pandas types to JSON serializable types"""
    import numpy as np
    import pandas as pd
    
    if isinstance(obj, (np.integer, np.int64, np.int32)):
        return int(obj)
    elif isinstance(obj, (np.floating, np.float64, np.float32)):
        if pd.isna(obj) or np.isinf(obj):
            return None
        return float(obj)
    elif isinstance(obj, (np.bool_, bool)):
        return bool(obj)
    elif isinstance(obj, np.ndarray):
        return obj.tolist()
    elif isinstance(obj, pd.Series):
        return obj.tolist()
    elif isinstance(obj, dict):
        return {k: make_json_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_json_serializable(item) for item in obj]
    elif pd.isna(obj):
        return None
    else:
        return obj


def analyze_dataset_context(df):
    """Analyze dataset and provide context-aware suggestions"""
    try:
        client = Groq(api_key=os.environ.get("GROQ_API_KEY"))
        
        # Get dataset summary
        summary = {
            'shape': df.shape,
            'columns': list(df.columns),
            'dtypes': df.dtypes.to_dict(),
            'numeric_cols': list(df.select_dtypes(include=[np.number]).columns),
            'categorical_cols': list(df.select_dtypes(include=['object']).columns),
            'sample_data': df.head(3).to_dict('records'),
            'missing_values': df.isnull().sum().to_dict()
        }
        
        prompt = f"""
        Analyze this dataset and provide intelligent, specific insights:
        
        Dataset Summary:
        - Shape: {summary['shape']} (rows, columns)
        - Columns: {summary['columns']}
        - Data types: {summary['dtypes']}
        - Numeric columns: {summary['numeric_cols']}
        - Categorical columns: {summary['categorical_cols']}
        - Sample data: {summary['sample_data']}
        - Missing values: {summary['missing_values']}
        
        Based on the actual column names and data types, provide:
        1. What domain/business context this data represents (be specific, not generic)
        2. 5 SPECIFIC SQL queries using actual column names that would reveal valuable business insights
        3. 3 SELECTIVE priority visualizations that make sense for this exact dataset (avoid overwhelming charts)
        4. Focus on the most important and insightful visualizations, not everything
        
        For visualizations, prefer:
        - Bar charts for top N categories (not all categories)
        - Scatter plots with meaningful relationships
        - Histograms for key numeric distributions
        - Avoid time series if no clear time column exists
        
        Format response as JSON:
        {{
            "context": "Specific description of what this data represents",
            "domain": "specific domain (e.g., 'e-commerce', 'healthcare', 'finance')",
            "suggested_queries": [
                "SELECT specific_column, COUNT(*) FROM data GROUP BY specific_column ORDER BY COUNT(*) DESC LIMIT 10",
                "SELECT AVG(specific_numeric_column) FROM data WHERE specific_condition"
            ],
            "priority_visualizations": [
                {{"type": "bar", "description": "Top 10 [specific column] distribution", "reason": "Shows most significant patterns"}},
                {{"type": "scatter", "description": "[specific columns] correlation analysis", "reason": "Reveals key relationships"}},
                {{"type": "histogram", "description": "[specific numeric column] distribution", "reason": "Shows data spread for key metric"}}
            ]
        }}
        """
        
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.1-8b-instant",
            temperature=0.3,
            max_tokens=1000
        )
        
        import json
        analysis = json.loads(chat_completion.choices[0].message.content)
        return analysis
        
    except Exception as e:
        print(f"Analysis error: {e}")
        # Create smarter fallback based on actual data
        sample_numeric = [col for col in df.select_dtypes(include=[np.number]).columns[:2]]
        sample_categorical = [col for col in df.select_dtypes(include=['object']).columns[:2]]
        
        fallback_queries = []
        fallback_queries.append(f"SELECT COUNT(*) as total_rows FROM data")
        
        if sample_categorical:
            fallback_queries.append(f"SELECT \"{sample_categorical[0]}\", COUNT(*) as count FROM data GROUP BY \"{sample_categorical[0]}\" ORDER BY count DESC LIMIT 10")
        
        if sample_numeric:
            fallback_queries.append(f"SELECT AVG(\"{sample_numeric[0]}\") as average, MIN(\"{sample_numeric[0]}\") as minimum, MAX(\"{sample_numeric[0]}\") as maximum FROM data")
        
        if len(sample_numeric) >= 2:
            fallback_queries.append(f"SELECT \"{sample_numeric[0]}\", \"{sample_numeric[1]}\" FROM data WHERE \"{sample_numeric[0]}\" IS NOT NULL AND \"{sample_numeric[1]}\" IS NOT NULL LIMIT 100")
        
        fallback_queries.append("SELECT * FROM data LIMIT 20")
        
        return {
            "context": f"Dataset with {df.shape[0]} rows and {df.shape[1]} columns containing {len(sample_numeric)} numeric and {len(sample_categorical)} text columns",
            "domain": "data-analysis", 
            "suggested_queries": fallback_queries[:5],
            "priority_visualizations": [
                {"type": "bar", "description": f"Top 10 {sample_categorical[0] if sample_categorical else 'categories'} distribution", "reason": "Shows most significant data patterns"},
                {"type": "scatter", "description": f"Relationship between {sample_numeric[0] if len(sample_numeric) > 0 else 'first'} and {sample_numeric[1] if len(sample_numeric) > 1 else 'second'} variables", "reason": "Reveals correlations between key metrics"} if len(sample_numeric) >= 2 else {"type": "histogram", "description": f"Distribution of {sample_numeric[0] if sample_numeric else 'numeric values'}", "reason": "Shows data spread"},
                {"type": "histogram", "description": f"Frequency distribution of {sample_numeric[0] if sample_numeric else 'values'}", "reason": "Shows data distribution and outliers"}
            ][:3]  # Ensure only 3 visualizations
        }


def generate_curious_insights(df):
    """
    The Curiosity Engine: Automatically analyzes a DataFrame to discover interesting patterns
    and generates fascinating, ready-to-ask questions for the user.
    
    Args:
        df: pandas DataFrame to analyze
        
    Returns:
        list: A list of 3-5 fascinating, insight-driven questions
    """
    try:
        print(f"🔍 Curiosity Engine: Starting analysis of dataset with {df.shape[0]} rows and {df.shape[1]} columns")
        
        # === AUTOMATED EDA ===
        findings = []
        
        # 1. Statistical Summary
        print("📊 Generating statistical summary...")
        numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
        
        if numeric_cols:
            stats_summary = df[numeric_cols].describe()
            findings.append(f"Statistical Summary: Dataset has {len(numeric_cols)} numeric columns. ")
            findings.append(f"Key statistics - Mean values range from {stats_summary.loc['mean'].min():.2f} to {stats_summary.loc['mean'].max():.2f}. ")
            findings.append(f"Standard deviations range from {stats_summary.loc['std'].min():.2f} to {stats_summary.loc['std'].max():.2f}. ")
        
        # 2. Top Correlations
        print("🔗 Finding correlations...")
        if len(numeric_cols) >= 2:
            corr_matrix = df[numeric_cols].corr()
            # Get upper triangle of correlation matrix
            upper_triangle = np.triu(np.ones(corr_matrix.shape), k=1).astype(bool)
            corr_pairs = corr_matrix.where(upper_triangle).stack().sort_values(ascending=False)
            
            # Top 5 positive correlations
            top_positive = corr_pairs.head(5)
            if len(top_positive) > 0:
                findings.append(f"\nSTRONG POSITIVE CORRELATIONS: ")
                for (col1, col2), corr_val in top_positive.items():
                    if abs(corr_val) > 0.3:  # Only report meaningful correlations
                        findings.append(f"- {col1} and {col2} have a correlation of {corr_val:.2f}. ")
            
            # Top 5 negative correlations
            top_negative = corr_pairs.tail(5)
            if len(top_negative) > 0:
                findings.append(f"\nSTRONG NEGATIVE CORRELATIONS: ")
                for (col1, col2), corr_val in top_negative.items():
                    if abs(corr_val) > 0.3:
                        findings.append(f"- {col1} and {col2} have a negative correlation of {corr_val:.2f}. ")
        
        # 3. Categorical Distribution Skew
        print("📈 Analyzing categorical distributions...")
        if categorical_cols:
            findings.append(f"\nCATEGORICAL PATTERNS: ")
            for col in categorical_cols[:5]:  # Limit to first 5 categorical columns
                value_counts = df[col].value_counts(normalize=True)
                if len(value_counts) > 0:
                    top_percentage = value_counts.iloc[0] * 100
                    if top_percentage > 80:
                        findings.append(f"- Column '{col}' is heavily skewed: '{value_counts.index[0]}' represents {top_percentage:.1f}% of all values. ")
                    elif len(value_counts) > 5:
                        findings.append(f"- Column '{col}' has {len(value_counts)} unique values, with top value '{value_counts.index[0]}' at {top_percentage:.1f}%. ")
        
        # 4. Outlier Detection using IQR
        print("🎯 Detecting outliers...")
        if numeric_cols:
            findings.append(f"\nOUTLIER DETECTION: ")
            for col in numeric_cols[:5]:  # Check first 5 numeric columns
                Q1 = df[col].quantile(0.25)
                Q3 = df[col].quantile(0.75)
                IQR = Q3 - Q1
                lower_bound = Q1 - 1.5 * IQR
                upper_bound = Q3 + 1.5 * IQR
                outliers = df[(df[col] < lower_bound) | (df[col] > upper_bound)]
                
                if len(outliers) > 0:
                    outlier_percentage = (len(outliers) / len(df)) * 100
                    findings.append(f"- Column '{col}' has {len(outliers)} outliers ({outlier_percentage:.1f}% of data), "
                                  f"ranging from {outliers[col].min():.2f} to {outliers[col].max():.2f}. ")
        
        # 5. Time-Series Patterns (Bonus)
        print("⏰ Checking for time-series patterns...")
        datetime_cols = df.select_dtypes(include=['datetime64']).columns.tolist()
        
        # Also check for columns that might be dates but stored as strings
        for col in categorical_cols[:10]:
            if any(keyword in col.lower() for keyword in ['date', 'time', 'year', 'month', 'day']):
                try:
                    pd.to_datetime(df[col].head(10), errors='raise')
                    datetime_cols.append(col)
                except:
                    pass
        
        if datetime_cols:
            findings.append(f"\nTIME-SERIES PATTERNS: ")
            for col in datetime_cols[:3]:
                try:
                    date_col = pd.to_datetime(df[col]) if df[col].dtype == 'object' else df[col]
                    date_range = (date_col.max() - date_col.min()).days
                    findings.append(f"- Column '{col}' spans {date_range} days, from {date_col.min()} to {date_col.max()}. ")
                except:
                    pass
        
        # Consolidate findings into comprehensive summary
        statistical_summary = "".join(findings)
        print(f"✅ Statistical analysis complete. Summary length: {len(statistical_summary)} characters")
        
        # === GENERATE CURIOUS QUESTIONS ===
        print("🤖 Calling Groq API to generate fascinating questions...")
        
        prompt = f"""You are a world-class data scientist and business intelligence expert. Analyze this dataset summary and generate CONCISE, ACTION-ORIENTED insights.

{statistical_summary}

Generate 3-5 SHORT, POWERFUL questions that reveal hidden insights. Each question should be:
- CONCISE (max 15 words)
- SPECIFIC (use actual column names)
- ACTIONABLE (lead to SQL queries)
- IMPACTFUL (business value)

FORMATTING RULES:
1. Use **bold** for key metrics/column names
2. Use specific numbers and percentages
3. Frame as discovery: "What drives...", "How does...", "Which factors..."
4. Focus on relationships, trends, and anomalies
5. NO fluff or unnecessary words

Return EXACTLY 3-5 questions in this format:
1. What drives **[key column]** performance in **[segment]**?
2. How does **[column1]** impact **[column2]** (correlation: X.XX)?
3. Which **[category]** shows the highest **[metric]** growth?

PRIORITIZE: Correlations > Outliers > Trends > Distributions

Keep questions SHORT, BOLD, and SCANNABLE!"""

        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,  # Higher temperature for more creative questions
            max_tokens=800
        )
        
        response_text = response.choices[0].message.content.strip()
        print(f"🎯 Received response from Groq API")
        
        # Parse the questions from the response
        questions = []
        lines = response_text.split('\n')
        
        for line in lines:
            line = line.strip()
            # Match lines starting with numbers (1., 2., etc.)
            if line and (line[0].isdigit() or line.startswith('-') or line.startswith('•')):
                # Remove numbering and clean up
                question = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
                if question and len(question) > 20:  # Ensure it's a substantial question
                    questions.append(question)
        
        # If parsing failed, try to split by double newlines or return raw lines
        if len(questions) == 0:
            questions = [q.strip() for q in response_text.split('\n\n') if q.strip() and len(q.strip()) > 20][:5]
        
        print(f"✅ Generated {len(questions)} curious questions")
        for i, q in enumerate(questions, 1):
            print(f"   {i}. {q[:100]}...")
        
        return questions[:5]  # Return maximum 5 questions
        
    except Exception as e:
        print(f"❌ Error in Curiosity Engine: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Return fallback questions if analysis fails
        return [
            "What are the key trends and patterns in this dataset?",
            "Which factors have the strongest relationships with each other?",
            "Are there any unusual outliers or anomalies worth investigating?",
            "What insights could help drive better business decisions?"
        ]


def auto_visualize_dataset(request):
    """Generate automatic visualizations for the dataset and save each to session"""
    if request.method == 'POST':
        # Set processing flag to prevent cleanup during operation
        request.session['auto_viz_processing'] = True
        request.session.modified = True
        
        # Use the active dataset from multi-dataset system
        path_in_bucket, display_name = get_active_dataset_path(request)
        
        print(f"DEBUG - Auto-visualize called with active dataset: {display_name}")
        print(f"DEBUG - Active dataset path: {path_in_bucket}")
        print(f"DEBUG - Session keys: {list(request.session.keys())}")
        
        if not path_in_bucket:
            request.session['auto_viz_processing'] = False
            request.session.modified = True
            return JsonResponse({'error': 'No dataset found in session. Please upload a dataset first.'}, status=400)
        
        try:
            print(f"DEBUG - Attempting to download from Supabase: {path_in_bucket}")
            # Load dataset with better error handling
            try:
                res = supabase.storage.from_(SUPABASE_BUCKET).download(path_in_bucket)
                print(f"DEBUG - Downloaded {len(res)} bytes from Supabase")
            except Exception as supabase_error:
                print(f"ERROR - Supabase download failed: {supabase_error}")
                request.session['auto_viz_processing'] = False
                request.session.modified = True
                return JsonResponse({
                    'error': f'Dataset file not found in storage. The file may have been cleaned up. Please re-upload your dataset. Error: {str(supabase_error)}'
                }, status=404)
            
            df = pd.read_csv(io.BytesIO(res))
            print(f"DEBUG - Loaded dataframe with {len(df)} rows, {len(df.columns)} columns")
            
            # Validate dataset size for visualization
            if len(df) > 10000:
                # Sample large datasets for analysis
                df_sample = df.sample(n=5000, random_state=42)
                print(f"Large dataset detected ({len(df)} rows). Using sample of 5000 rows for analysis.")
            else:
                df_sample = df
            
            # Get analysis
            analysis = analyze_dataset_context(df_sample)
            
            # Limit to maximum 3 visualizations to avoid overwhelming
            priority_visualizations = analysis.get('priority_visualizations', [])[:3]
            
            # Generate the priority visualizations
            charts = []
            successful_charts = 0
            max_charts = 3
            
            for viz in priority_visualizations:
                if successful_charts >= max_charts:
                    break
                    
                try:
                    chart_data = generate_auto_chart(df_sample, viz['type'], viz['description'])
                    if chart_data:
                        # Save the visualization to session for report builder
                        save_visualization(
                            request,
                            f"{viz['type'].title()} Chart: {viz['description'][:50]}",
                            chart_data['image'],
                            chart_data['explanation']
                        )
                        charts.append({
                            'type': viz['type'],
                            'description': viz['description'],
                            'reason': viz['reason'],
                            'image': chart_data['image'],
                            'explanation': chart_data['explanation'],
                            'chart_title': f"{viz['type'].title()} Chart"
                        })
                        successful_charts += 1
                        print(f"Successfully generated {viz['type']} chart")
                    else:
                        print(f"Failed to generate {viz['type']} chart - no data returned")
                except Exception as e:
                    print(f"Failed to generate {viz['type']} chart: {e}")
                    continue
            
            if not charts:
                request.session['auto_viz_processing'] = False
                request.session.modified = True
                return JsonResponse({
                    'error': 'Could not generate any visualizations for this dataset. Please try manual chart creation.',
                    'analysis': analysis
                })
            
            # Clear processing flag on success
            request.session['auto_viz_processing'] = False
            request.session.modified = True
            
            return JsonResponse({
                'success': True,
                'analysis': analysis,
                'charts': charts,
                'message': f'Generated {len(charts)} selective visualizations focusing on key insights.'
            })
            
        except Exception as e:
            # Clear processing flag on error
            request.session['auto_viz_processing'] = False
            request.session.modified = True
            
            print(f"ERROR - Auto-visualization failed: {str(e)}")
            print(f"ERROR - Exception type: {type(e)}")
            import traceback
            print(f"ERROR - Traceback: {traceback.format_exc()}")
            return JsonResponse({
                'error': f'Auto-visualization failed: {str(e)}. This may be due to session cleanup. Please try re-uploading your dataset.'
            }, status=500)
    
    return JsonResponse({'error': 'Only POST method allowed'}, status=405)

def generate_auto_chart(df, chart_type, description):
    """Generate automatic charts using pure Python (Matplotlib/Seaborn) - NO Chart.js"""
    try:
        import matplotlib
        matplotlib.use('Agg')  # Use non-interactive backend
        import matplotlib.pyplot as plt
        import seaborn as sns
        import base64
        from io import BytesIO
        import warnings
        warnings.filterwarnings('ignore')
        
        # Set style for professional charts
        plt.style.use('dark_background')
        sns.set_palette("viridis")  # Changed from husl to viridis for better auto-charts
        
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        categorical_cols = df.select_dtypes(include=['object']).columns
        
        # Limit data size for better performance and readability
        max_categories = 10
        max_data_points = 100
        
        # Create figure
        plt.figure(figsize=(10, 6))
        
        if chart_type == 'bar' and len(categorical_cols) > 0 and len(numeric_cols) > 0:
            # Bar chart of categorical vs numeric
            cat_col = categorical_cols[0]
            num_col = numeric_cols[0]
            
            # Get top N categories to avoid overcrowding
            top_categories = df[cat_col].value_counts().head(max_categories)
            filtered_df = df[df[cat_col].isin(top_categories.index)]
            
            # Group and aggregate data
            grouped = filtered_df.groupby(cat_col)[num_col].mean().round(2)
            
            # Create bar chart
            bars = plt.bar(range(len(grouped)), grouped.values, color='#7059f2', alpha=0.8)
            plt.xticks(range(len(grouped)), grouped.index, rotation=45, ha='right')
            plt.xlabel(cat_col)
            plt.ylabel(f'Average {num_col}')
            plt.title(f'Top {len(grouped)} {cat_col} by Average {num_col}', fontsize=14, fontweight='bold')
            
            # Add value labels on bars
            for bar, value in zip(bars, grouped.values):
                plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01*max(grouped.values),
                        f'{value:.1f}', ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            explanation = f'Bar chart showing top {len(grouped)} categories of {cat_col} by average {num_col}. Limited to most significant categories for clarity.'
            
        elif chart_type == 'scatter' and len(numeric_cols) >= 2:
            # Scatter plot of two numeric columns
            x_col, y_col = numeric_cols[0], numeric_cols[1]
            
            # Sample data if too large
            sample_df = df.sample(n=min(max_data_points, len(df))).copy()
            
            # Remove outliers for better visualization
            for col in [x_col, y_col]:
                Q1 = sample_df[col].quantile(0.25)
                Q3 = sample_df[col].quantile(0.75)
                IQR = Q3 - Q1
                lower_bound = Q1 - 1.5 * IQR
                upper_bound = Q3 + 1.5 * IQR
                sample_df = sample_df[(sample_df[col] >= lower_bound) & (sample_df[col] <= upper_bound)]
            
            # Create scatter plot
            plt.scatter(sample_df[x_col], sample_df[y_col], color='#7059f2', alpha=0.6, s=30)
            plt.xlabel(x_col)
            plt.ylabel(y_col)
            plt.title(f'Scatter Plot: {x_col} vs {y_col}', fontsize=14, fontweight='bold')
            plt.grid(True, alpha=0.3)
            plt.tight_layout()
            
            explanation = f'Scatter plot showing relationship between {x_col} and {y_col}. Data sampled to {len(sample_df)} points and outliers removed for better visualization.'
            
        elif chart_type == 'histogram' and len(numeric_cols) > 0:
            # Histogram of first numeric column
            num_col = numeric_cols[0]
            
            # Remove outliers
            Q1 = df[num_col].quantile(0.25)
            Q3 = df[num_col].quantile(0.75)
            IQR = Q3 - Q1
            lower_bound = Q1 - 1.5 * IQR
            upper_bound = Q3 + 1.5 * IQR
            filtered_data = df[(df[num_col] >= lower_bound) & (df[num_col] <= upper_bound)][num_col]
            
            # Create histogram
            plt.hist(filtered_data.dropna(), bins=min(30, len(filtered_data)//10), 
                    color='#7059f2', alpha=0.7, edgecolor='white', linewidth=0.5)
            plt.xlabel(num_col)
            plt.ylabel('Frequency')
            plt.title(f'Distribution of {num_col}', fontsize=14, fontweight='bold')
            plt.grid(True, alpha=0.3)
            plt.tight_layout()
            
            explanation = f'Histogram showing distribution of {num_col}. Outliers removed for better visualization.'
            
        else:
            # Fallback: Simple bar chart of column value counts
            if len(categorical_cols) > 0:
                cat_col = categorical_cols[0]
                value_counts = df[cat_col].value_counts().head(max_categories)
                
                plt.bar(range(len(value_counts)), value_counts.values, color='#7059f2', alpha=0.8)
                plt.xticks(range(len(value_counts)), value_counts.index, rotation=45, ha='right')
                plt.xlabel(cat_col)
                plt.ylabel('Count')
                plt.title(f'Distribution of {cat_col}', fontsize=14, fontweight='bold')
                plt.tight_layout()
                
                explanation = f'Bar chart showing distribution of {cat_col} values.'
            else:
                # If no categorical columns, show numeric distribution
                if len(numeric_cols) > 0:
                    num_col = numeric_cols[0]
                    plt.hist(df[num_col].dropna(), bins=30, color='#7059f2', alpha=0.7)
                    plt.xlabel(num_col)
                    plt.ylabel('Frequency')
                    plt.title(f'Distribution of {num_col}', fontsize=14, fontweight='bold')
                    plt.tight_layout()
                    explanation = f'Histogram showing distribution of {num_col}.'
                else:
                    plt.close()
                    return None
        
        # Convert to base64 image
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        plt.close('all')  # Clean up
        
        return {
            'image': f'data:image/png;base64,{image_base64}',
            'explanation': explanation
        }
        
    except Exception as e:
        plt.close('all')  # Cleanup on error
        print(f"Error generating Python chart: {e}")
        return None


def generate_chart_image(chart_type, x_data, y_data, title):
    """Generate a simple base64 encoded chart image using matplotlib"""
    try:
        import matplotlib.pyplot as plt
        import base64
        from io import BytesIO
        
        plt.style.use('default')
        fig, ax = plt.subplots(figsize=(8, 6))
        
        if chart_type == 'bar':
            ax.bar(range(len(x_data)), y_data, color='#7059f2')
            ax.set_xticks(range(len(x_data)))
            ax.set_xticklabels(x_data, rotation=45, ha='right')
        elif chart_type == 'scatter':
            ax.scatter(x_data, y_data, color='#7059f2', alpha=0.6)
        elif chart_type == 'histogram':
            ax.bar(range(len(x_data)), y_data, color='#7059f2')
            ax.set_xticks(range(len(x_data)))
            ax.set_xticklabels([f'{x:.1f}' for x in x_data], rotation=45, ha='right')
        elif chart_type == 'pie':
            ax.pie(y_data, labels=x_data, autopct='%1.1f%%', startangle=90)
        
        ax.set_title(title, fontsize=14, fontweight='bold')
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', dpi=100, bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.read()).decode()
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating chart image: {e}")
        return ""



def dataset_preview(request, dataset_name):
    # Use active dataset from multi-dataset system
    path_in_bucket, display_name = get_active_dataset_path(request)
    if not display_name:
        display_name = dataset_name  # fallback
    print(f"DEBUG: Preview active dataset: {display_name}")  # DEBUG
    print(f"DEBUG: Preview dataset path: {path_in_bucket}")  # DEBUG

    if not path_in_bucket:
        messages.error(request, "No dataset in session.")
        return redirect("index")

    try:
        res = supabase.storage.from_(SUPABASE_BUCKET).download(path_in_bucket)
        df = pd.read_csv(io.BytesIO(res))
        print(f"✅ Loaded dataset from Supabase: {dataset_name}")  # DEBUG

        # Generate comprehensive data profile
        data_profile = generate_data_profile(df)

        # Analyze dataset context and provide suggestions
        analysis = analyze_dataset_context(df)
        
        # 🔍 CURIOSITY ENGINE: Generate fascinating, auto-discovered questions
        print("🚀 Activating Curiosity Engine...")
        curious_questions = generate_curious_insights(df)
        print(f"✨ Curiosity Engine generated {len(curious_questions)} questions")
        
        # Save dataset schema to session for report builder
        save_dataset_schema(request, display_name, data_profile)

    except Exception as e:
        df = pd.DataFrame()
        analysis = None
        data_profile = None
        curious_questions = []
        print(f"❌ Failed to load dataset from Supabase: {e}")  # DEBUG
        messages.error(
            request,
            "Dataset not found. It may have been cleaned up or your session expired. Please upload the dataset again."
        )
        return redirect("index")

    context = {
        "dataset": {
            "name": dataset_name,            # internal Supabase name
            "display_name": display_name     # clean original filename
        },
        "columns": df.columns.tolist() if not df.empty else [],
        "preview_rows": df.head(5).values.tolist() if not df.empty else [],
        "analysis": analysis,
        "data_profile": data_profile,
        "curious_questions": curious_questions,  # 🔍 Curiosity Engine questions
    }

    return render(request, "dataset_preview.html", context)



def full_dataset(request):
    # Use the active dataset from multi-dataset system
    path_in_bucket, display_name = get_active_dataset_path(request)
    print(f"DEBUG: full_dataset active dataset: {display_name}")  # DEBUG
    print(f"DEBUG: full_dataset path: {path_in_bucket}")  # DEBUG

    if not path_in_bucket:
        return JsonResponse({"error": "No dataset in session."}, status=400)

    try:
        res = supabase.storage.from_(SUPABASE_BUCKET).download(path_in_bucket)
        df = pd.read_csv(io.BytesIO(res))

        # Get query parameters for pagination
        limit = int(request.GET.get('limit', 1000))
        offset = int(request.GET.get('offset', 0))
        
        # Apply pagination
        total_rows = len(df)
        df_page = df.iloc[offset:offset+limit]

        def make_json_safe(x):
            if isinstance(x, (np.integer, int)):
                return int(x)
            if isinstance(x, (np.floating, float)):
                if pd.isna(x) or np.isinf(x):
                    return None
                return float(x)
            if pd.isna(x):
                return None
            if isinstance(x, (np.bool_, bool)):
                return bool(x)
            return str(x)

        safe_rows = [{col: make_json_safe(row[col]) for col in df_page.columns} for _, row in df_page.iterrows()]

        return JsonResponse({
            "columns": df.columns.tolist(),
            "rows": safe_rows,
            "total_rows": total_rows,
            "offset": offset,
            "limit": limit,
            "has_more": offset + limit < total_rows
        })

    except Exception as e:
        print(f"❌ full_dataset failed: {e}")  # DEBUG
        return JsonResponse({"error": f"Failed to read dataset: {str(e)}"}, status=500)


# Helper function to clean data for JSON serialization
def sanitize_for_json(obj):
    """Convert numpy/pandas types to JSON-serializable Python types"""
    import numpy as np
    
    if isinstance(obj, (np.integer, np.int64, np.int32)):
        return int(obj)
    elif isinstance(obj, (np.floating, np.float64, np.float32)):
        if np.isnan(obj) or np.isinf(obj):
            return None
        return float(obj)
    elif isinstance(obj, np.ndarray):
        return [sanitize_for_json(item) for item in obj.tolist()]
    elif isinstance(obj, list):
        return [sanitize_for_json(item) for item in obj]
    elif isinstance(obj, dict):
        return {key: sanitize_for_json(val) for key, val in obj.items()}
    elif isinstance(obj, tuple):
        return tuple(sanitize_for_json(item) for item in obj)
    elif pd.isna(obj):
        return None
    else:
        return obj


# Helper functions for SQL query generation
def make_sql_prompt(dataset_name: str, columns: list[dict], user_input: str) -> str:
    """Generate prompt for AI to convert natural language to SQL"""
    obj = {
        "role": "system",
        "task": "Translate English to exactly one SQLite SELECT query.",
        "dialect": "sqlite",
        "table": {
            "name": dataset_name,
            "columns": columns  # [{"name":"col","type":"int64"}, ...]
        },
        "rules": [
            "Respond with JSON ONLY. No prose, no code fences.",
            "Return exactly one key: sql.",
            "The value must be exactly one SELECT statement ending with a semicolon.",
            "Use ONLY the provided table.",
            f'Always double-quote the table name in FROM: FROM "{dataset_name}".',
            "Use only columns that exist.",
            "Do not use JOINs, other tables, PRAGMA, or DDL.",
            "Numeric literals unquoted; strings single-quoted and escaped.",
            "No comments, no multiple statements."
        ],
        "input": {"english": user_input},
        "output_schema": {
            "type": "object",
            "required": ["sql"],
            "additionalProperties": False,
            "properties": {"sql": {"type": "string"}}
        },
        "respond_with": "json_only"
    }
    return json.dumps(obj, ensure_ascii=False)


def extract_top_level_json(text: str) -> dict:
    """Extract JSON from AI response, handling various formats"""
    # First try direct parse
    try:
        return json.loads(text)
    except Exception:
        pass
    # Fallback: extract first balanced {...} block
    depth = 0
    start = -1
    for i, ch in enumerate(text):
        if ch == '{':
            if depth == 0:
                start = i
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0 and start != -1:
                block = text[start:i+1]
                try:
                    return json.loads(block)
                except Exception:
                    # keep scanning in case there is another block later
                    start = -1
    raise ValueError("Model did not return valid JSON.")


def run_query(request, dataset_name):
    import re

    if request.method != "POST":
        return JsonResponse({"error": "POST request required"}, status=400)

    # Handle both JSON and form data
    if request.content_type == 'application/json':
        try:
            data = json.loads(request.body)
            user_input = data.get("query")
        except json.JSONDecodeError:
            return JsonResponse({"error": "Invalid JSON"}, status=400)
    else:
        user_input = request.POST.get("query")
    
    # Use active dataset from multi-dataset system
    path_in_bucket, display_name = get_active_dataset_path(request)
    print(f"DEBUG: run_query active dataset: {display_name}")
    print(f"DEBUG: run_query dataset path: {path_in_bucket}")

    if not path_in_bucket:
        return JsonResponse({
            "query": user_input or "", 
            "result": [], 
            "error": "No dataset found in session. Please upload a new dataset to continue."
        })

    # Load CSV
    try:
        res = supabase.storage.from_(SUPABASE_BUCKET).download(path_in_bucket)
        df = pd.read_csv(io.BytesIO(res))
        print(f"✅ Dataset loaded for query: {dataset_name}")
    except Exception as e:
        print(f"❌ Failed to load dataset: {e}")
        error_msg = "Dataset file not found. This may happen if:"
        error_msg += "<br>• The file was automatically cleaned up"
        error_msg += "<br>• Your session has expired"
        error_msg += "<br>• There was a storage issue"
        error_msg += "<br><br>Please <a href='/' style='color: #4f46e5; text-decoration: underline;'>upload your dataset again</a> to continue."
        
        return JsonResponse({
            "query": user_input or "", 
            "result": [], 
            "error": error_msg
        })

    # Build in-memory SQLite table
    conn = sqlite3.connect(":memory:")
    try:
        # Let pandas create the table with proper quoting
        df.to_sql(dataset_name, conn, index=False, if_exists="replace")
    except Exception as e:
        conn.close()
        return JsonResponse({"query": user_input or "", "result": [], "error": f"Failed to stage table: {str(e)}"})

    # Column metadata for the prompt
    columns = [{"name": c, "type": str(t)} for c, t in zip(df.columns, df.dtypes)]

    # JSON prompt (safe)
    prompt_json = make_sql_prompt(dataset_name, columns, user_input)

    try:
        # DEBUG: Check if client and API key are available
        print(f"DEBUG: client object type: {type(client)}")
        print(f"DEBUG: GROQ_API_KEY exists: {bool(os.getenv('GROQ_API_KEY'))}")
        print(f"DEBUG: GROQ_API_KEY first 20 chars: {os.getenv('GROQ_API_KEY')[:20] if os.getenv('GROQ_API_KEY') else 'NOT FOUND'}")
        print(f"DEBUG: Calling Groq API with model: llama-3.1-8b-instant")
        
        # Call model — keep temp 0; send the JSON as the sole message content
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt_json}],
            temperature=0,
            # If your client supports it, uncomment to force JSON:
            # response_format={"type": "json_object"},
        )
        print(f"DEBUG: API call successful!")
        raw = response.choices[0].message.content.strip()
        obj = extract_top_level_json(raw)

        if list(obj.keys()) != ["sql"]:
            raise ValueError("Bad shape: expected exactly one key 'sql'.")

        sql_query = obj["sql"].strip()

        # Guards
        if not re.match(r'(?is)^\s*select\b', sql_query):
            return JsonResponse({"query": sql_query, "result": [], "error": "Only SELECT queries are allowed."})

        # Auto-fix semicolon if missing
        if not sql_query.endswith(";"):
            sql_query = sql_query + ";"

        # Check for multiple statements (semicolon in the middle)
        if re.search(r";\s*\S", sql_query):
            return JsonResponse({"query": sql_query, "result": [], "error": "Only single SELECT statements are allowed."})

        # Simple table name check - just ensure the dataset name appears somewhere in the query
        # This is more flexible than strict regex matching
        if dataset_name not in sql_query:
            return JsonResponse({"query": sql_query, "result": [], "error": f'Query must reference the dataset "{dataset_name}".'})

        print(f"DEBUG: Table name validation passed for: {dataset_name}")
        
        # Optional: add LIMIT 1000 if none present (prevents giant tables)
        auto_limit_added = False
        if not re.search(r'(?is)\blimit\s+\d+\b', sql_query):
            sql_query = sql_query[:-1] + " LIMIT 1000;"
            auto_limit_added = True

        print(f"DEBUG: Generated SQL: {sql_query}")

        # Execute
        try:
            result_df = pd.read_sql_query(sql_query, conn)
        except Exception as e:
            cols = df.columns.tolist()
            return JsonResponse({
                "query": sql_query,
                "result": [],
                "error": f"Query failed: {str(e)}. Available columns: {', '.join(cols)}",
                "auto_limit_added": auto_limit_added
            })

        if result_df.empty:
            return JsonResponse({"query": sql_query, "result": [], "auto_limit_added": auto_limit_added})

        # Convert DataFrame to list of dictionaries for JSON response
        def make_json_safe(x):
            if isinstance(x, (np.integer, int)):
                return int(x)
            if isinstance(x, (np.floating, float)):
                if pd.isna(x) or np.isinf(x):
                    return None
                return float(x)
            if pd.isna(x):
                return None
            if isinstance(x, (np.bool_, bool)):
                return bool(x)
            return str(x)

        result_data = []
        for _, row in result_df.iterrows():
            row_dict = {}
            for col in result_df.columns:
                row_dict[col] = make_json_safe(row[col])
            result_data.append(row_dict)

        # Add visualization suggestions
        visualization_suggestion = suggest_visualization_type(result_df, sql_query)

        # ---- Save to session for report builder ----
        # Save the user query as a chat message
        save_chat_message(request, "User", user_input)
        
        # Save the AI response (SQL + results summary)
        ai_response = f"Generated SQL Query:\n```sql\n{sql_query}\n```\n\nFound {len(result_data)} results."
        save_chat_message(request, "AI Assistant", ai_response)
        
        # Save the detailed query result
        save_query_result(
            request,
            user_input,
            sql_query,
            result_data,
            visualization_suggestion
        )
        
        # Also save as a standalone table for easy report inclusion
        save_query_table(
            request,
            f"Query Results: {user_input[:50]}...",
            result_data,
            user_input,
            sql_query
        )
        # -------------------------------------------

        return JsonResponse({
            "query": sql_query, 
            "result": result_data,
            "visualization_suggestion": visualization_suggestion,
            "auto_limit_added": auto_limit_added,
            "data_summary": {
                "rows": len(result_data),
                "columns": list(result_df.columns),
                "numeric_columns": list(result_df.select_dtypes(include=[np.number]).columns),
                "categorical_columns": list(result_df.select_dtypes(include=['object']).columns)
            }
        })

    except Exception as e:
        # DEBUG: Print full error details
        import traceback
        print(f"DEBUG ERROR: {type(e).__name__}: {str(e)}")
        print(f"DEBUG ERROR TRACEBACK:\n{traceback.format_exc()}")
        return JsonResponse({"query": user_input or "", "result": [], "error": f"Failed to generate/parse SQL: {str(e)}"})

    finally:
        conn.close()



def save_report_blocks(request):
    """Save edited report blocks back to session"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            blocks = data.get('blocks', [])
            request.session['report_blocks'] = blocks
            return JsonResponse({'status': 'success'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Invalid method'})


def save_report(request, format):
    """Enhanced report generation with complete chat history and figures"""
    try:
        if request.method == 'POST':
            data = json.loads(request.body)
            blocks = data.get("blocks", [])
            filename = data.get("filename")  # Optional custom filename from client (without extension)
            
            print(f"[SAVE_REPORT DEBUG] Received POST request for {format} format")
            print(f"[SAVE_REPORT DEBUG] Number of blocks: {len(blocks)}")
            
            for i, block in enumerate(blocks):
                print(f"[SAVE_REPORT DEBUG] Block {i}: {block.get('type', 'unknown')} - {list(block.keys())}")
                if block.get('has_table'):
                    print(f"[SAVE_REPORT DEBUG] Block {i} has table with {len(block.get('result_data', []))} rows")
            
            # If no blocks provided, get from session
            if not blocks:
                blocks = request.session.get('report_blocks', [])
                print(f"[SAVE_REPORT DEBUG] Using session blocks: {len(blocks)}")
            
        else:
            # GET request - use session data
            blocks = request.session.get('report_blocks', [])
            print(f"[SAVE_REPORT DEBUG] GET request, using session blocks: {len(blocks)}")
        
        if format == "pdf":
            return generate_comprehensive_pdf(blocks, filename)
        elif format == "pptx":
            return generate_comprehensive_pptx(blocks)
        elif format in ["jpg", "png"]:
            return generate_report_image(blocks, format, filename)
            
    except Exception as e:
        print(f"[SAVE_REPORT DEBUG] Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'error': f'Report generation failed: {str(e)}'}, status=500)
    
    return HttpResponse("Invalid format", status=400)


def generate_comprehensive_pdf(blocks, filename=None):
    """Generate a comprehensive PDF with chat history, figures, and styling"""
    print(f"[PDF DEBUG] Generating PDF with {len(blocks)} blocks")
    
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        import tempfile
        import requests
        import re
        from html import unescape
    except ImportError as e:
        # Fallback to simple PDF generation if reportlab is not available
        print(f"ReportLab not available: {e}")
        return generate_simple_pdf(blocks)
    
    try:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, 
                              rightMargin=50, leftMargin=50, 
                              topMargin=50, bottomMargin=50)
        
        # Create styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.HexColor('#2E86AB'),
            alignment=1  # Center alignment
        )
        
        user_style = ParagraphStyle(
            'UserMessage',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=10,
            leftIndent=20,
            textColor=colors.HexColor('#1E3A8A'),
            backColor=colors.HexColor('#EBF8FF')
        )
        
        ai_style = ParagraphStyle(
            'AIMessage',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=10,
            leftIndent=20,
            textColor=colors.HexColor('#059669'),
            backColor=colors.HexColor('#ECFDF5')
        )
        
        code_style = ParagraphStyle(
            'CodeBlock',
            parent=styles['Code'],
            fontSize=10,
            leftIndent=30,
            backColor=colors.HexColor('#F3F4F6'),
            textColor=colors.HexColor('#374151')
        )
        
        story = []
        
        # Add professional title page
        story.append(Spacer(1, 100))
        story.append(Paragraph("Data Analysis Report", title_style))
        story.append(Spacer(1, 12))
        
        # Add subtitle with better formatting
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=styles['Normal'],
            fontSize=12,
            textColor=colors.HexColor('#666666'),
            alignment=1  # Center
        )
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", subtitle_style))
        story.append(Spacer(1, 50))
        
        # Add a separator line
        from reportlab.platypus import HRFlowable
        story.append(HRFlowable(width="80%", thickness=2, color=colors.HexColor('#4F46E5'), spaceBefore=10, spaceAfter=30))
        
        # Process blocks
        for i, block in enumerate(blocks):
            block_type = block.get('type', '')
            print(f"[PDF DEBUG] Processing block {i}: {block_type}")
            print(f"[PDF DEBUG] Block data: {block}")
            
            if block_type == 'chat_message':
                sender = block.get('sender', 'User')
                content = block.get('content', '')
                has_table = block.get('has_table', False)
                result_data = block.get('result_data', [])
                
                print(f"[PDF DEBUG] Chat message - Sender: {sender}, Has table: {has_table}, Result data length: {len(result_data) if result_data else 0}")
                
                if sender.lower() == 'user':
                    story.append(Paragraph(f"<b>User Question:</b>", user_style))
                    # Clean HTML content for user messages
                    clean_content = re.sub('<[^<]+?>', '', content) if content else ''
                    story.append(Paragraph(clean_content, user_style))
                else:
                    story.append(Paragraph(f"<b>AI Response:</b>", ai_style))
                    
                    # Check if this AI message contains a table
                    if has_table and result_data:
                        print(f"[PDF DEBUG] Processing AI message with table: {len(result_data)} rows")
                        
                        # Extract text content without HTML tags for the main response
                        content_text = block.get('content_text', content)
                        if content_text:
                            clean_content = re.sub('<[^<]+?>', '', content_text)
                            story.append(Paragraph(clean_content, ai_style))
                        
                        story.append(Spacer(1, 10))
                        story.append(Paragraph(f"<b>Data Table ({len(result_data)} rows):</b>", styles['Heading4']))
                        
                        # Create table for PDF
                        headers = list(result_data[0].keys()) if result_data else []
                        table_data = [headers]  # Header row
                        
                        print(f"[PDF DEBUG] Table headers: {headers}")
                        
                        # Add data rows (limit to first 15 rows for PDF space)
                        for row_idx, row in enumerate(result_data[:15]):
                            row_data = [str(row.get(col, '')) for col in headers]
                            table_data.append(row_data)
                            if row_idx < 3:  # Log first few rows
                                print(f"[PDF DEBUG] Row {row_idx}: {row_data}")
                        
                        if len(table_data) > 1:  # If we have data beyond headers
                            try:
                                table = Table(table_data)
                                table.setStyle(TableStyle([
                                    ('BACKGROUND', (0, 0), (-1, 0), colors.black),
                                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                    ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                                    ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                                ]))
                                story.append(table)
                                story.append(Spacer(1, 15))
                                print(f"[PDF DEBUG] Table added successfully")
                                
                                # Add note if data was truncated
                                if len(result_data) > 15:
                                    story.append(Paragraph(f"<i>Note: Table shows first 15 rows of {len(result_data)} total results.</i>", styles['Normal']))
                                    story.append(Spacer(1, 10))
                            except Exception as table_error:
                                print(f"[PDF DEBUG] Error creating table in chat message: {table_error}")
                                story.append(Paragraph(f"Table data: {len(result_data)} rows", styles['Normal']))
                    else:
                        # Regular AI message without table
                        print(f"[PDF DEBUG] Regular AI message without table")
                        clean_content = re.sub('<[^<]+?>', '', content) if content else ''
                        story.append(Paragraph(clean_content, ai_style))
                
                story.append(Spacer(1, 15))
                
            elif block_type == 'query':
                user_input = block.get('user_input', '')
                sql_query = block.get('sql_query', '')
                result_data = block.get('result_data', [])
                
                story.append(Paragraph(f"<b>Query Request:</b>", user_style))
                story.append(Paragraph(user_input, user_style))
                story.append(Spacer(1, 10))
                
                story.append(Paragraph(f"<b>Generated SQL:</b>", code_style))
                story.append(Paragraph(f"<font name='Courier'>{sql_query}</font>", code_style))
                story.append(Spacer(1, 15))
                
                # Add table data if available
                if result_data and len(result_data) > 0:
                    story.append(Paragraph(f"<b>Results ({len(result_data)} rows):</b>", styles['Heading4']))
                    
                    # Create table for PDF
                    headers = list(result_data[0].keys()) if result_data else []
                    table_data = [headers]  # Header row
                    
                    # Add data rows (limit to first 10 rows for PDF space)
                    for i, row in enumerate(result_data[:10]):
                        row_data = [str(row.get(col, '')) for col in headers]
                        table_data.append(row_data)
                    
                    if len(table_data) > 1:  # If we have data beyond headers
                        try:
                            table = Table(table_data)
                            table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, 0), 10),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
                                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                                ('FONTSIZE', (0, 1), (-1, -1), 8),
                                ('GRID', (0, 0), (-1, -1), 1, colors.black)
                            ]))
                            story.append(table)
                            story.append(Spacer(1, 20))
                            
                            # Add note if data was truncated
                            if len(result_data) > 10:
                                story.append(Paragraph(f"<i>Note: Table shows first 10 rows of {len(result_data)} total results.</i>", styles['Normal']))
                                story.append(Spacer(1, 10))
                        except Exception as table_error:
                            print(f"Error creating table: {table_error}")
                            story.append(Paragraph(f"Table data: {len(result_data)} rows returned", styles['Normal']))
                            story.append(Spacer(1, 15))
                
            elif block_type == 'visualization':
                chart_title = block.get('chart_title', 'Chart')
                image_url = block.get('image_url', '')
                explanation = block.get('explanation', '')
                pdf_w_in = block.get('pdf_width_inch')
                pdf_h_in = block.get('pdf_height_inch')

                story.append(Paragraph(f"<b>Visualization: {chart_title}</b>", styles['Heading3']))

                # Embed image from base64 or URL
                if image_url:
                    try:
                        tmp_file_path = None
                        if image_url.startswith('data:'):
                            # Base64 data URL
                            header, encoded = image_url.split(',', 1)
                            image_data = base64.b64decode(encoded)
                            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                                tmp_file.write(image_data)
                                tmp_file_path = tmp_file.name
                        elif image_url.startswith('http'):
                            # Remote HTTP/HTTPS image
                            response = requests.get(image_url, timeout=10)
                            if response.status_code == 200:
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                                    tmp_file.write(response.content)
                                    tmp_file_path = tmp_file.name
                            else:
                                story.append(Paragraph(f"[Image request failed: status {response.status_code}]", styles['Normal']))
                        else:
                            # Local/relative path (static/media)
                            candidate = None
                            if os.path.isabs(image_url) and os.path.exists(image_url):
                                candidate = image_url
                            elif image_url.startswith('/'):
                                candidate = os.path.join(settings.BASE_DIR, image_url.lstrip('/'))
                            else:
                                candidate = os.path.join(settings.BASE_DIR, image_url.replace('/', os.sep))
                            if candidate and os.path.exists(candidate):
                                tmp_file_path = candidate
                        # If we created a temp file, add image to PDF
                        if tmp_file_path:
                            img = Image(tmp_file_path)
                            # Apply custom sizing if provided
                            try:
                                if pdf_w_in or pdf_h_in:
                                    # Determine original image size
                                    orig_w = getattr(img, 'imageWidth', None)
                                    orig_h = getattr(img, 'imageHeight', None)
                                    if pdf_w_in and not pdf_h_in and orig_w and orig_h:
                                        # Width only, keep aspect ratio
                                        scale = (float(pdf_w_in) * inch) / float(orig_w)
                                        img.drawWidth = float(pdf_w_in) * inch
                                        img.drawHeight = float(orig_h) * scale
                                    elif pdf_h_in and not pdf_w_in and orig_w and orig_h:
                                        # Height only, keep aspect ratio
                                        scale = (float(pdf_h_in) * inch) / float(orig_h)
                                        img.drawHeight = float(pdf_h_in) * inch
                                        img.drawWidth = float(orig_w) * scale
                                    elif pdf_w_in and pdf_h_in:
                                        # Explicit width and height
                                        img.drawWidth = float(pdf_w_in) * inch
                                        img.drawHeight = float(pdf_h_in) * inch
                                else:
                                    # Default max bounds
                                    img._restrictSize(6*inch, 4*inch)
                            except Exception:
                                # Fallback fixed size
                                img.drawWidth = 6*inch
                                img.drawHeight = 4*inch
                            story.append(img)
                            import os
                            try:
                                # Only delete if it's a temp file we created (not a permanent local file)
                                if tmp_file_path and tmp_file_path.startswith(tempfile.gettempdir()):
                                    os.unlink(tmp_file_path)
                            except Exception:
                                pass
                    except Exception as e:
                        story.append(Paragraph(f"[Image could not be loaded: {str(e)}]", styles['Normal']))

                if explanation:
                    story.append(Spacer(1, 10))
                    story.append(Paragraph(explanation, styles['Normal']))

                story.append(Spacer(1, 20))
                
            elif block_type == 'dataset_schema':
                dataset_name = block.get('dataset_name', 'Dataset')
                shape = block.get('shape', (0, 0))
                columns = block.get('columns', [])
                data_health = block.get('data_health', {})
                
                story.append(Paragraph(f"<b>Dataset Schema: {dataset_name}</b>", styles['Heading3']))
                story.append(Paragraph(f"Shape: {shape[0]} rows × {shape[1]} columns", styles['Normal']))
                
                # Data health summary
                if data_health:
                    health_text = f"Missing data: {data_health.get('missing_percentage', 0):.1f}% | "
                    health_text += f"Duplicate rows: {data_health.get('duplicate_rows', 0)}"
                    story.append(Paragraph(health_text, styles['Normal']))
                
                # Column information table
                if columns:
                    col_headers = ['Column', 'Type', 'Missing %', 'Unique Count']
                    col_data = [col_headers]
                    
                    for col in columns[:10]:  # Limit to first 10 columns
                        row = [
                            col.get('name', ''),
                            col.get('dtype', ''),
                            f"{col.get('missing_percentage', 0):.1f}%",
                            str(col.get('unique_count', ''))
                        ]
                        col_data.append(row)
                    
                    try:
                        schema_table = Table(col_data)
                        schema_table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
                            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                            ('FONTSIZE', (0, 1), (-1, -1), 8),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(schema_table)
                    except Exception as e:
                        story.append(Paragraph(f"Schema information: {len(columns)} columns", styles['Normal']))
                
                story.append(Spacer(1, 20))
                
            elif block_type == 'query_table':
                table_title = block.get('title', 'Query Results')
                table_data = block.get('data', [])
                user_input = block.get('user_input', '')
                sql_query = block.get('sql_query', '')
                
                story.append(Paragraph(f"<b>{table_title}</b>", styles['Heading3']))
                
                if user_input:
                    story.append(Paragraph(f"<b>Query:</b> {user_input}", user_style))
                if sql_query:
                    story.append(Paragraph(f"<b>SQL:</b> <font name='Courier'>{sql_query}</font>", code_style))
                    story.append(Spacer(1, 10))
                
                # Add table data
                if table_data and len(table_data) > 0:
                    headers = list(table_data[0].keys())
                    table_rows = [headers]  # Header row
                    
                    # Add data rows (limit to first 10 rows for PDF space)
                    for i, row in enumerate(table_data[:10]):
                        row_data = [str(row.get(col, '')) for col in headers]
                        table_rows.append(row_data)
                    
                    try:
                        result_table = Table(table_rows)
                        result_table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
                            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                            ('FONTSIZE', (0, 1), (-1, -1), 8),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(result_table)
                        
                        # Add note if data was truncated
                        if len(table_data) > 10:
                            story.append(Spacer(1, 10))
                            story.append(Paragraph(f"<i>Note: Table shows first 10 rows of {len(table_data)} total results.</i>", styles['Normal']))
                    except Exception as table_error:
                        print(f"Error creating query table: {table_error}")
                        story.append(Paragraph(f"Query returned {len(table_data)} rows", styles['Normal']))
                
                story.append(Spacer(1, 20))
                
            elif block_type == 'text':
                content = block.get('content', '')
                # Clean HTML tags and decode entities
                clean_content = re.sub('<[^<]+?>', '', content) if content else ''
                clean_content = unescape(clean_content)
                if clean_content.strip():
                    story.append(Paragraph(clean_content, styles['Normal']))
                    story.append(Spacer(1, 15))
                    
            elif block_type == 'title':
                title_text = block.get('title', block.get('content', ''))
                # Clean HTML and get plain text
                clean_title = re.sub('<[^<]+?>', '', title_text) if title_text else ''
                clean_title = unescape(clean_title)
                if clean_title.strip():
                    # Create custom title style for h1 headings
                    heading_style = ParagraphStyle(
                        'ReportTitle',
                        parent=styles['Heading2'],
                        fontSize=18,
                        spaceAfter=20,
                        spaceBefore=10,
                        textColor=colors.HexColor('#2E86AB'),
                        fontName='Helvetica-Bold'
                    )
                    story.append(Paragraph(clean_title, heading_style))
                    story.append(Spacer(1, 10))
                    
            elif block_type == 'sql':
                sql_query = block.get('sql_query', block.get('content', ''))
                if sql_query:
                    story.append(Paragraph("<b>SQL Query:</b>", styles['Heading4']))
                    story.append(Paragraph(f"<font name='Courier' size=9>{sql_query}</font>", code_style))
                    story.append(Spacer(1, 15))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        # Determine filename
        safe_name = (filename or 'data_analysis_report').strip()
        if not safe_name:
            safe_name = 'data_analysis_report'
        # Remove illegal filename characters
        safe_name = ''.join(ch for ch in safe_name if ch not in '\\/:*?"<>|').strip()
        response = HttpResponse(buffer, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename={safe_name}.pdf'
        return response
        
    except Exception as e:
        print(f"Advanced PDF generation failed: {e}")
        return generate_simple_pdf(blocks, filename)


def generate_simple_pdf(blocks, filename=None):
    """Simple PDF generation fallback when reportlab is not available"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
    except ImportError:
        # Final fallback - return HTML
        return generate_html_fallback(blocks)
    
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    y_position = height - 50
    p.setFont("Helvetica-Bold", 16)
    p.drawString(50, y_position, "Data Analysis Report")
    y_position -= 30
    
    p.setFont("Helvetica", 10)
    p.drawString(50, y_position, f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
    y_position -= 40
    
    p.setFont("Helvetica", 10)
    
    for block in blocks:
        if y_position < 100:  # Start new page
            p.showPage()
            y_position = height - 50
        
        block_type = block.get('type', '')
        
        if block_type == 'chat_message':
            sender = block.get('sender', 'User')
            content = block.get('content', '')[:100]  # Truncate for simple PDF
            
            p.setFont("Helvetica-Bold", 10)
            p.drawString(50, y_position, f"{sender}:")
            y_position -= 15
            
            p.setFont("Helvetica", 9)
            p.drawString(70, y_position, content)
            y_position -= 25
            
        elif block_type == 'visualization':
            chart_title = block.get('chart_title', 'Chart')
            explanation = block.get('explanation', '')[:80]  # Truncate
            
            p.setFont("Helvetica-Bold", 10)
            p.drawString(50, y_position, f"Chart: {chart_title}")
            y_position -= 15
            
            if explanation:
                p.setFont("Helvetica", 9)
                p.drawString(70, y_position, explanation)
                y_position -= 25
    
    p.save()
    buffer.seek(0)
    
    safe_name = (filename or 'data_analysis_report').strip()
    if not safe_name:
        safe_name = 'data_analysis_report'
    safe_name = ''.join(ch for ch in safe_name if ch not in '\\/:*?"<>|').strip()
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename={safe_name}.pdf'
    return response


def generate_html_fallback(blocks):
    """HTML fallback when PDF generation is not possible"""
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Data Analysis Report</title>
        <style>
            body {{ font-family: Arial, sans-serif; padding: 20px; max-width: 1000px; margin: 0 auto; }}
            .user-message {{ background: #e3f2fd; padding: 10px; margin: 10px 0; border-radius: 5px; }}
            .ai-message {{ background: #e8f5e8; padding: 10px; margin: 10px 0; border-radius: 5px; }}
            .query {{ background: #f5f5f5; padding: 15px; margin: 15px 0; border-radius: 5px; border-left: 4px solid #4F46E5; }}
            .schema {{ background: #fff3cd; padding: 15px; margin: 15px 0; border-radius: 5px; border-left: 4px solid #ffc107; }}
            .query_table {{ background: #d1ecf1; padding: 15px; margin: 15px 0; border-radius: 5px; border-left: 4px solid #17a2b8; }}
            .chart {{ text-align: center; margin: 20px 0; }}
            .chart img {{ max-width: 100%; }}
            table {{ border-collapse: collapse; width: 100%; margin: 10px 0; font-size: 12px; }}
            th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
            th {{ background-color: #4F46E5; color: white; font-weight: bold; }}
            tr:nth-child(even) {{ background-color: #f9f9f9; }}
            code {{ background: #f0f0f0; padding: 2px 4px; border-radius: 3px; font-family: monospace; }}
            h1 {{ color: #333; text-align: center; }}
            h3, h4 {{ color: #4F46E5; }}
        </style>
    </head>
    <body>
        <h1>Data Analysis Report</h1>
        <p><strong>Generated on:</strong> {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
    """
    
    for block in blocks:
        block_type = block.get('type', '')
        
        if block_type == 'chat_message':
            sender = block.get('sender', 'User')
            content = block.get('content', '')
            css_class = 'user-message' if sender.lower() == 'user' else 'ai-message'
            html_content += f'<div class="{css_class}"><strong>{sender}:</strong> {content}</div>'
            
        elif block_type == 'query':
            user_input = block.get('user_input', '')
            sql_query = block.get('sql_query', '')
            result_data = block.get('result_data', [])
            
            html_content += f'<div class="query"><h3>Query Request</h3>'
            html_content += f'<p><strong>Question:</strong> {user_input}</p>'
            html_content += f'<p><strong>SQL:</strong> <code>{sql_query}</code></p>'
            
            # Add table if data exists
            if result_data and len(result_data) > 0:
                html_content += f'<h4>Results ({len(result_data)} rows):</h4>'
                html_content += '<table border="1" style="border-collapse: collapse; width: 100%; margin: 10px 0;">'
                
                # Headers
                headers = list(result_data[0].keys())
                html_content += '<tr style="background-color: #4F46E5; color: white;">'
                for header in headers:
                    html_content += f'<th style="padding: 8px;">{header}</th>'
                html_content += '</tr>'
                
                # Data rows (limit to first 10)
                for row in result_data[:10]:
                    html_content += '<tr>'
                    for header in headers:
                        html_content += f'<td style="padding: 8px;">{row.get(header, "")}</td>'
                    html_content += '</tr>'
                
                html_content += '</table>'
                
                if len(result_data) > 10:
                    html_content += f'<p><i>Note: Showing first 10 rows of {len(result_data)} total results.</i></p>'
            
            html_content += '</div>'
            
        elif block_type == 'dataset_schema':
            dataset_name = block.get('dataset_name', 'Dataset')
            shape = block.get('shape', (0, 0))
            columns = block.get('columns', [])
            data_health = block.get('data_health', {})
            
            html_content += f'<div class="schema"><h3>Dataset Schema: {dataset_name}</h3>'
            html_content += f'<p><strong>Shape:</strong> {shape[0]} rows × {shape[1]} columns</p>'
            
            if data_health:
                html_content += f'<p><strong>Data Health:</strong> '
                html_content += f'Missing data: {data_health.get("missing_percentage", 0):.1f}% | '
                html_content += f'Duplicate rows: {data_health.get("duplicate_rows", 0)}</p>'
            
            if columns:
                html_content += '<h4>Column Information:</h4>'
                html_content += '<table border="1" style="border-collapse: collapse; width: 100%; margin: 10px 0;">'
                html_content += '<tr style="background-color: #4F46E5; color: white;">'
                html_content += '<th style="padding: 8px;">Column</th>'
                html_content += '<th style="padding: 8px;">Type</th>'
                html_content += '<th style="padding: 8px;">Missing %</th>'
                html_content += '<th style="padding: 8px;">Unique Count</th>'
                html_content += '</tr>'
                
                for col in columns[:10]:  # Limit to first 10 columns
                    html_content += '<tr>'
                    html_content += f'<td style="padding: 8px;">{col.get("name", "")}</td>'
                    html_content += f'<td style="padding: 8px;">{col.get("dtype", "")}</td>'
                    html_content += f'<td style="padding: 8px;">{col.get("missing_percentage", 0):.1f}%</td>'
                    html_content += f'<td style="padding: 8px;">{col.get("unique_count", "")}</td>'
                    html_content += '</tr>'
                
                html_content += '</table>'
                
                if len(columns) > 10:
                    html_content += f'<p><i>Note: Showing first 10 columns of {len(columns)} total columns.</i></p>'
            
            html_content += '</div>'
            
        elif block_type == 'query_table':
            table_title = block.get('title', 'Query Results')
            table_data = block.get('data', [])
            user_input = block.get('user_input', '')
            sql_query = block.get('sql_query', '')
            
            html_content += f'<div class="query_table"><h3>{table_title}</h3>'
            
            if user_input:
                html_content += f'<p><strong>Query:</strong> {user_input}</p>'
            if sql_query:
                html_content += f'<p><strong>SQL:</strong> <code>{sql_query}</code></p>'
            
            if table_data and len(table_data) > 0:
                html_content += '<table border="1" style="border-collapse: collapse; width: 100%; margin: 10px 0;">'
                
                # Headers
                headers = list(table_data[0].keys())
                html_content += '<tr style="background-color: #4F46E5; color: white;">'
                for header in headers:
                    html_content += f'<th style="padding: 8px;">{header}</th>'
                html_content += '</tr>'
                
                # Data rows (limit to first 10)
                for row in table_data[:10]:
                    html_content += '<tr>'
                    for header in headers:
                        html_content += f'<td style="padding: 8px;">{row.get(header, "")}</td>'
                    html_content += '</tr>'
                
                html_content += '</table>'
                
                if len(table_data) > 10:
                    html_content += f'<p><i>Note: Showing first 10 rows of {len(table_data)} total results.</i></p>'
            
            html_content += '</div>'
            
        elif block_type == 'visualization':
            chart_title = block.get('chart_title', 'Chart')
            image_url = block.get('image_url', '')
            explanation = block.get('explanation', '')
            
            html_content += f'<div class="chart"><h3>{chart_title}</h3>'
            if image_url:
                html_content += f'<img src="{image_url}" alt="{chart_title}">'
            if explanation:
                html_content += f'<p>{explanation}</p>'
            html_content += '</div>'
    
    html_content += """
    </body>
    </html>
    """
    
    response = HttpResponse(html_content, content_type='text/html')
    response['Content-Disposition'] = 'attachment; filename=report.html'
    return response
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    import tempfile
    import requests
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, 
                          rightMargin=50, leftMargin=50, 
                          topMargin=50, bottomMargin=50)
    
    # Create styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.HexColor('#2E86AB'),
        alignment=1  # Center alignment
    )
    
    user_style = ParagraphStyle(
        'UserMessage',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=10,
        leftIndent=20,
        textColor=colors.HexColor('#1E3A8A'),
        backColor=colors.HexColor('#EBF8FF')
    )
    
    ai_style = ParagraphStyle(
        'AIMessage',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=10,
        leftIndent=20,
        textColor=colors.HexColor('#059669'),
        backColor=colors.HexColor('#ECFDF5')
    )
    
    code_style = ParagraphStyle(
        'CodeBlock',
        parent=styles['Code'],
        fontSize=10,
        leftIndent=30,
        backColor=colors.HexColor('#F3F4F6'),
        textColor=colors.HexColor('#374151')
    )
    
    story = []
    
    # Add title
    story.append(Paragraph("Data Analysis Report", title_style))
    story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", styles['Normal']))
    story.append(Spacer(1, 30))
    
    # Process blocks
    for i, block in enumerate(blocks):
        block_type = block.get('type', '')
        
        if block_type == 'chat_message':
            sender = block.get('sender', 'User')
            content = block.get('content', '')
            
            if sender.lower() == 'user':
                story.append(Paragraph(f"<b>User Question:</b>", user_style))
                story.append(Paragraph(content, user_style))
            else:
                story.append(Paragraph(f"<b>AI Response:</b>", ai_style))
                story.append(Paragraph(content, ai_style))
            
            story.append(Spacer(1, 15))
            
        elif block_type == 'query':
            user_input = block.get('user_input', '')
            sql_query = block.get('sql_query', '')
            
            story.append(Paragraph(f"<b>Query Request:</b>", user_style))
            story.append(Paragraph(user_input, user_style))
            story.append(Spacer(1, 10))
            
            story.append(Paragraph(f"<b>Generated SQL:</b>", code_style))
            story.append(Paragraph(f"<font name='Courier'>{sql_query}</font>", code_style))
            story.append(Spacer(1, 15))
            
        elif block_type == 'visualization':
            chart_title = block.get('chart_title', 'Chart')
            image_url = block.get('image_url', '')
            explanation = block.get('explanation', '')
            
            story.append(Paragraph(f"<b>Visualization: {chart_title}</b>", styles['Heading3']))
            
            # Download and add image
            if image_url and (image_url.startswith('http') or image_url.startswith('data:')):
                try:
                    if image_url.startswith('data:'):
                        # Handle base64 images
                        header, encoded = image_url.split(',', 1)
                        image_data = base64.b64decode(encoded)
                        
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                            tmp_file.write(image_data)
                            tmp_file_path = tmp_file.name
                        
                        img = Image(tmp_file_path, width=5*inch, height=3*inch)
                        story.append(img)
                        
                        # Clean up temp file
                        import os
                        os.unlink(tmp_file_path)
                    else:
                        # Handle URL images
                        response = requests.get(image_url)
                        if response.status_code == 200:
                            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                                tmp_file.write(response.content)
                                tmp_file_path = tmp_file.name
                            
                            img = Image(tmp_file_path, width=5*inch, height=3*inch)
                            story.append(img)
                            
                            # Clean up temp file
                            import os
                            os.unlink(tmp_file_path)
                except Exception as e:
                    story.append(Paragraph(f"[Image could not be loaded: {str(e)}]", styles['Normal']))
            
            if explanation:
                story.append(Spacer(1, 10))
                story.append(Paragraph(explanation, styles['Normal']))
            
            story.append(Spacer(1, 20))
            
        elif block_type == 'text':
            content = block.get('content', '')
            story.append(Paragraph(content, styles['Normal']))
            story.append(Spacer(1, 15))
            
        elif block_type == 'table':
            # Handle table data if present
            table_data = block.get('data', [])
            if table_data:
                # Convert to reportlab table
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F46E5')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8FAFC')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
                story.append(Spacer(1, 20))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename=data_analysis_report.pdf'
    return response


def generate_comprehensive_pptx(blocks):
    """Generate comprehensive PowerPoint with chat history and figures"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    import tempfile
    import requests
    
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Data Analysis Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    
    current_slide = None
    slide_content = []
    
    for block in blocks:
        block_type = block.get('type', '')
        
        if block_type == 'visualization':
            # Create new slide for each visualization
            slide_layout = prs.slide_layouts[5]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = block.get('chart_title', 'Visualization')
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Add image
            image_url = block.get('image_url', '')
            if image_url:
                try:
                    from PIL import Image as PILImage
                    if image_url.startswith('data:'):
                        # Handle base64 images
                        header, encoded = image_url.split(',', 1)
                        image_data = base64.b64decode(encoded)
                        
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                            tmp_file.write(image_data)
                            tmp_file_path = tmp_file.name
                        # Preserve aspect ratio
                        wpx, hpx = PILImage.open(tmp_file_path).size
                        target_w = Inches(8)
                        target_h = Inches(8 * hpx / float(wpx))
                        slide.shapes.add_picture(tmp_file_path, Inches(1), Inches(2), target_w, target_h)
                        
                        # Clean up
                        import os
                        os.unlink(tmp_file_path)
                    else:
                        tmp_file_path = None
                        if image_url.startswith('http'):
                            # Handle URL images
                            response = requests.get(image_url, timeout=10)
                            if response.status_code == 200:
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                                    tmp_file.write(response.content)
                                    tmp_file_path = tmp_file.name
                        else:
                            # Local/relative path
                            candidate = None
                            if os.path.isabs(image_url) and os.path.exists(image_url):
                                candidate = image_url
                            elif image_url.startswith('/'):
                                candidate = os.path.join(settings.BASE_DIR, image_url.lstrip('/'))
                            else:
                                candidate = os.path.join(settings.BASE_DIR, image_url.replace('/', os.sep))
                            if candidate and os.path.exists(candidate):
                                tmp_file_path = candidate
                        if tmp_file_path:
                            wpx, hpx = PILImage.open(tmp_file_path).size
                            target_w = Inches(8)
                            target_h = Inches(8 * hpx / float(wpx))
                            slide.shapes.add_picture(tmp_file_path, Inches(1), Inches(2), target_w, target_h)
                            # Clean up temp file if created
                            import os
                            if tmp_file_path.startswith(tempfile.gettempdir()):
                                os.unlink(tmp_file_path)
                except Exception as e:
                    # Add error text
                    error_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
                    error_frame = error_shape.text_frame
                    error_frame.text = f"Image could not be loaded: {str(e)}"
            
            # Add explanation
            explanation = block.get('explanation', '')
            if explanation:
                exp_shape = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(1))
                exp_frame = exp_shape.text_frame
                exp_frame.text = explanation
        
        elif block_type in ['chat_message', 'query', 'text']:
            # Accumulate text content for text slides
            if block_type == 'chat_message':
                sender = block.get('sender', 'User')
                content = block.get('content', '')
                has_table = block.get('has_table', False)
                result_data = block.get('result_data', [])
                
                slide_content.append(f"{sender}: {content}")
                
                # If chat message has table data, add it
                if has_table and result_data:
                    slide_content.append(f"Data Table ({len(result_data)} rows):")
                    # Add first few rows as sample
                    if result_data:
                        headers = list(result_data[0].keys())
                        slide_content.append(f"Columns: {', '.join(headers)}")
                        for i, row in enumerate(result_data[:3]):  # Show first 3 rows
                            row_text = ', '.join([f"{k}: {v}" for k, v in row.items()])
                            slide_content.append(f"Row {i+1}: {row_text}")
                        if len(result_data) > 3:
                            slide_content.append(f"... and {len(result_data) - 3} more rows")
                            
            elif block_type == 'query':
                user_input = block.get('user_input', '')
                sql_query = block.get('sql_query', '')
                result_data = block.get('result_data', [])
                has_table = block.get('has_table', False)
                
                slide_content.append(f"Query: {user_input}")
                slide_content.append(f"SQL: {sql_query}")
                
                if has_table and result_data:
                    slide_content.append(f"Results: {len(result_data)} rows")
                    # Add sample data
                    if result_data:
                        headers = list(result_data[0].keys())
                        slide_content.append(f"Columns: {', '.join(headers)}")
                        for i, row in enumerate(result_data[:3]):  # Show first 3 rows
                            row_text = ', '.join([f"{k}: {v}" for k, v in row.items()])
                            slide_content.append(f"Row {i+1}: {row_text}")
                        if len(result_data) > 3:
                            slide_content.append(f"... and {len(result_data) - 3} more rows")
                elif result_data:
                    slide_content.append(f"Results: {len(result_data)} rows returned")
                    
            elif block_type == 'text':
                content = block.get('content', '')
                slide_content.append(content)
    
    # Add accumulated text content to a summary slide
    if slide_content:
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = "Analysis Summary"
        content_text = '\n\n'.join(slide_content[:10])  # Limit content
        content_placeholder.text = content_text
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    response = HttpResponse(buffer.getvalue(), 
                          content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response["Content-Disposition"] = "attachment; filename=data_analysis_report.pptx"
    return response


def clear_session_data_endpoint(request):
    """Endpoint to manually clear session data"""
    if request.method == 'POST':
        clear_session_data(request)
        # Also clear dataset-related session keys
        dataset_keys = ['datasets', 'active_dataset_index', 'uploaded_files', 'current_dataset_path']
        for key in dataset_keys:
            if key in request.session:
                del request.session[key]
        request.session.modified = True
        return JsonResponse({'status': 'success', 'message': 'All session data cleared. Please re-upload your files.'})
    return JsonResponse({'status': 'error', 'message': 'POST request required'})


def get_available_charts(request):
    """Get all available charts/visualizations from session"""
    try:
        visualizations = request.session.get('visualizations', [])
        charts = []
        
        for i, viz in enumerate(visualizations):
            charts.append({
                'id': i,
                'title': viz.get('chart_title', f'Chart {i+1}'),
                'image_url': viz.get('image_url', ''),
                'explanation': viz.get('explanation', ''),
                'timestamp': viz.get('timestamp', '')
            })
        
        return JsonResponse({'status': 'success', 'charts': charts})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)})


def get_chat_history(request):
    """Get chat history from session for report builder"""
    try:
        # Debug session contents
        print(f"DEBUG - All session keys: {list(request.session.keys())}")
        
        # Get chat history from session
        chat_history = request.session.get('chat_history', [])
        
        print(f"DEBUG - Chat history length: {len(chat_history)}")
        if chat_history:
            print(f"DEBUG - First chat item: {chat_history[0]}")
        
        # Format chat history for report builder
        formatted_history = []
        for chat in chat_history:
            # Handle different chat types
            if chat.get('type') == 'chat_message':
                formatted_history.append({
                    'sender': chat.get('sender', 'User'),
                    'content': chat.get('content', ''),
                    'timestamp': chat.get('timestamp', '')
                })
            elif chat.get('type') == 'query':
                # Add user query as a chat message
                formatted_history.append({
                    'sender': 'User',
                    'content': chat.get('user_input', ''),
                    'timestamp': chat.get('timestamp', '')
                })
                
                # Create comprehensive AI response with SQL and results
                sql_query = chat.get('sql_query', '')
                result_data = chat.get('result_data', [])
                
                # Build response content with SQL and table
                response_content = f"<strong>Generated SQL Query:</strong><br><pre><code>{sql_query}</code></pre>"
                
                # Add results table if we have data
                if result_data and len(result_data) > 0:
                    response_content += "<br><strong>Query Results:</strong><br>"
                    response_content += "<table style='border-collapse: collapse; width: 100%; margin-top: 10px;'>"
                    
                    # Table header
                    headers = list(result_data[0].keys())
                    response_content += "<tr style='background-color: #f0f0f0;'>"
                    for header in headers:
                        response_content += f"<th style='border: 1px solid #ddd; padding: 8px; text-align: left;'>{header}</th>"
                    response_content += "</tr>"
                    
                    # Table rows (limit to first 20 for display)
                    for i, row in enumerate(result_data[:20]):
                        response_content += "<tr>"
                        for header in headers:
                            value = row.get(header, '')
                            response_content += f"<td style='border: 1px solid #ddd; padding: 8px;'>{value}</td>"
                        response_content += "</tr>"
                    
                    response_content += "</table>"
                    
                    # Add note if there are more rows
                    if len(result_data) > 20:
                        response_content += f"<br><em>Showing first 20 of {len(result_data)} total rows</em>"
                
                # Add AI response about the query
                formatted_history.append({
                    'sender': 'AI Assistant',
                    'content': response_content,
                    'timestamp': chat.get('timestamp', '')
                })
        
        print(f"DEBUG - Formatted history length: {len(formatted_history)}")
        
        return JsonResponse({
            'success': True,
            'chat_history': formatted_history,
            'debug_info': {
                'session_keys': list(request.session.keys()),
                'raw_chat_history_length': len(chat_history),
                'formatted_history_length': len(formatted_history)
            }
        })
        
    except Exception as e:
        print(f"ERROR - get_chat_history failed: {str(e)}")
        import traceback
        print(f"ERROR - Traceback: {traceback.format_exc()}")
        return JsonResponse({
            'success': False,
            'error': str(e),
            'chat_history': []
        })


def generate_report_image(blocks, format_type, filename=None):
    """Export images from report blocks.
    - Collects all images from visualization and image blocks
    - Applies optional width/height (inches) if provided
    - Returns a single image when only one is present, or a ZIP of multiple images
    """
    from PIL import Image, ImageOps, ImageDraw, ImageFont
    import zipfile
    import tempfile
    import re
    import requests
    from django.utils.text import slugify

    images = []

    # Basic rendering helpers for non-image blocks
    CANVAS_W = 1280
    MARGIN = 40
    BG_COLOR = 'white'
    FG_COLOR = 'black'

    def get_font(size=16, bold=False):
        try:
            # Try a common font if present; fallback to default
            return ImageFont.truetype("arial.ttf", size)
        except Exception:
            return ImageFont.load_default()

    def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int):
        lines = []
        for paragraph in (text or '').split('\n'):
            words = paragraph.split(' ')
            cur = ''
            for w in words:
                trial = (cur + ' ' + w).strip()
                bbox = draw.textbbox((0,0), trial, font=font)
                if bbox[2] - bbox[0] <= max_width:
                    cur = trial
                else:
                    if cur:
                        lines.append(cur)
                    cur = w
            if cur:
                lines.append(cur)
            # paragraph break
            lines.append('')
        # Remove trailing empty line
        if lines and lines[-1] == '':
            lines.pop()
        return lines

    def render_text_block(title: str = None, text: str = '', subtitle: str = None):
        # Create a canvas tall enough for content (estimate, then grow if needed)
        img = Image.new('RGB', (CANVAS_W, 200), BG_COLOR)
        draw = ImageDraw.Draw(img)
        title_font = get_font(28)
        body_font = get_font(18)
        y = MARGIN
        max_text_w = CANVAS_W - 2*MARGIN
        # Title
        if title:
            draw.text((MARGIN, y), title, fill=FG_COLOR, font=title_font)
            y += 42
        if subtitle:
            draw.text((MARGIN, y), subtitle, fill=FG_COLOR, font=body_font)
            y += 32
        # Body
        lines = wrap_text(draw, text, body_font, max_text_w)
        for line in lines:
            if line == '':
                y += 12
                continue
            draw.text((MARGIN, y), line, fill=FG_COLOR, font=body_font)
            y += 26
        # Extend canvas height if needed
        final_h = max(y + MARGIN, 200)
        if final_h != img.size[1]:
            new_img = Image.new('RGB', (CANVAS_W, final_h), BG_COLOR)
            new_img.paste(img, (0,0))
            img = new_img
        return img

    def render_table(title: str, headers: list, rows: list):
        # Limit rows for image size
        headers = headers or []
        rows = rows or []
        max_rows = 30
        rows = rows[:max_rows]
        title_font = get_font(24)
        cell_font = get_font(16)
        draw_dummy = ImageDraw.Draw(Image.new('RGB', (10,10)))
        # Compute column widths evenly
        cols = len(headers)
        cols = max(cols, 1)
        table_w = CANVAS_W - 2*MARGIN
        col_w = table_w // cols
        row_h = 28
        header_h = 34
        table_h = header_h + len(rows)*row_h
        # Canvas height includes title and margins
        top_area = MARGIN + 40  # title space
        img_h = top_area + table_h + MARGIN
        img = Image.new('RGB', (CANVAS_W, img_h), BG_COLOR)
        draw = ImageDraw.Draw(img)
        y = MARGIN
        # Title
        if title:
            draw.text((MARGIN, y), title, fill=FG_COLOR, font=title_font)
        y += 40
        # Header background
        draw.rectangle([MARGIN, y, MARGIN+table_w, y+header_h], outline='black', fill='#eeeeee')
        # Header cells
        for ci in range(cols):
            x0 = MARGIN + ci*col_w
            x1 = x0 + col_w
            draw.rectangle([x0, y, x1, y+header_h], outline='black')
            head_text = str(headers[ci]) if ci < len(headers) else ''
            draw.text((x0+6, y+8), head_text, fill=FG_COLOR, font=get_font(16))
        y += header_h
        # Rows
        for r in rows:
            for ci in range(cols):
                x0 = MARGIN + ci*col_w
                x1 = x0 + col_w
                draw.rectangle([x0, y, x1, y+row_h], outline='#999999')
                cell_text = str(r.get(headers[ci], '')) if ci < len(headers) else ''
                draw.text((x0+6, y+6), cell_text, fill=FG_COLOR, font=cell_font)
            y += row_h
        return img

    def load_image_from_url(url: str):
        if url.startswith('data:'):
            try:
                header, encoded = url.split(',', 1)
                raw = base64.b64decode(encoded)
                return Image.open(BytesIO(raw))
            except Exception as e:
                print(f"[IMG EXPORT] Failed to decode base64 image: {e}")
                return None
        elif url.startswith('http'):
            try:
                r = requests.get(url, timeout=15)
                r.raise_for_status()
                return Image.open(BytesIO(r.content))
            except Exception as e:
                print(f"[IMG EXPORT] HTTP image fetch failed: {e}")
                return None
        else:
            # Local path resolution
            candidate = None
            if os.path.isabs(url) and os.path.exists(url):
                candidate = url
            elif url.startswith('/'):
                candidate = os.path.join(settings.BASE_DIR, url.lstrip('/'))
            else:
                candidate = os.path.join(settings.BASE_DIR, url.replace('/', os.sep))
            try:
                if candidate and os.path.exists(candidate):
                    return Image.open(candidate)
            except Exception as e:
                print(f"[IMG EXPORT] Local image open failed: {e}")
            return None

    DPI = 96.0

    for idx, block in enumerate(blocks):
        btype = block.get('type')
        if btype not in ('visualization', 'image'):
            # Render non-image blocks
            if btype == 'title':
                title_text = block.get('title') or block.get('content', '')
                img = render_text_block(title=title_text)
                safe_name = slugify((title_text or 'title')[:80]) or f"title-{idx+1:03d}"
                images.append((safe_name, img))
                continue
            if btype == 'text':
                content = block.get('content', '')
                # strip basic tags
                try:
                    import re
                    content = re.sub('<[^<]+?>', '', content)
                except Exception:
                    pass
                img = render_text_block(text=content)
                safe_name = f"text-{idx+1:03d}"
                images.append((safe_name, img))
                continue
            if btype == 'chat_message':
                sender = block.get('sender', 'Message')
                content = block.get('content', '')
                try:
                    import re
                    content = re.sub('<[^<]+?>', '', content)
                except Exception:
                    pass
                img = render_text_block(title=sender, text=content)
                safe_name = slugify(sender) or f"chat-{idx+1:03d}"
                images.append((safe_name, img))
                continue
            if btype == 'query':
                title = 'Query Results'
                data = block.get('result_data', [])
                headers = list(data[0].keys()) if data else []
                img = render_table(title, headers, data)
                safe_name = f"query-{idx+1:03d}"
                images.append((safe_name, img))
                continue
            if btype == 'query_table':
                title = block.get('title', 'Query Table')
                data = block.get('data', [])
                headers = list(data[0].keys()) if data else []
                img = render_table(title, headers, data)
                safe_name = slugify(title) or f"table-{idx+1:03d}"
                images.append((safe_name, img))
                continue
            # Skip other block types
            continue
        url = block.get('image_url')
        if not url:
            continue
        img = load_image_from_url(url)
        if img is None:
            continue

        # Apply sizing if provided
        w_in = block.get('pdf_width_inch')
        h_in = block.get('pdf_height_inch')
        try:
            if w_in or h_in:
                ow, oh = img.size
                if w_in and not h_in:
                    new_w = int(float(w_in) * DPI)
                    new_h = int(oh * (new_w / ow))
                    img = img.resize((max(1, new_w), max(1, new_h)), Image.LANCZOS)
                elif h_in and not w_in:
                    new_h = int(float(h_in) * DPI)
                    new_w = int(ow * (new_h / oh))
                    img = img.resize((max(1, new_w), max(1, new_h)), Image.LANCZOS)
                elif w_in and h_in:
                    new_w = int(float(w_in) * DPI)
                    new_h = int(float(h_in) * DPI)
                    img = img.resize((max(1, new_w), max(1, new_h)), Image.LANCZOS)
        except Exception as e:
            print(f"[IMG EXPORT] Resize failed: {e}")

        title = block.get('chart_title') or block.get('alt_text') or f"image_{idx+1}"
        safe_name = slugify(title) or f"image-{idx+1:03d}"
        images.append((safe_name, img))

    if not images:
        return JsonResponse({'error': 'No images found to export'}, status=400)

    # Single image: return directly
    if len(images) == 1:
        name, img = images[0]
        buf = BytesIO()
        if format_type == 'jpg':
            # Ensure RGB for JPEG and flatten transparency onto white
            if img.mode in ('RGBA', 'P'):
                bg = Image.new('RGB', img.size, 'white')
                try:
                    alpha = img.split()[3] if img.mode == 'RGBA' else None
                    bg.paste(img, mask=alpha)
                    img = bg
                except Exception:
                    img = img.convert('RGB')
            else:
                img = img.convert('RGB')
            img.save(buf, format='JPEG', quality=90)
            content_type = 'image/jpeg'
            ext = 'jpg'
        else:
            img.save(buf, format='PNG')
            content_type = 'image/png'
            ext = 'png'
        buf.seek(0)
        download_name = (filename or name) + f'.{ext}'
        resp = HttpResponse(buf.getvalue(), content_type=content_type)
        resp['Content-Disposition'] = f'attachment; filename="{download_name}"'
        return resp

    # Multiple images: zip them
    zip_buf = BytesIO()
    with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, (name, img) in enumerate(images, start=1):
            img_buf = BytesIO()
            if format_type == 'jpg':
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                img.save(img_buf, format='JPEG', quality=90)
                ext = 'jpg'
            else:
                img.save(img_buf, format='PNG')
                ext = 'png'
            img_buf.seek(0)
            filename_in_zip = f"{i:03d}_{name}.{ext}"
            zf.writestr(filename_in_zip, img_buf.getvalue())
    zip_buf.seek(0)
    download_name = (filename or 'report_images') + '.zip'
    response = HttpResponse(zip_buf.getvalue(), content_type='application/zip')
    response['Content-Disposition'] = f'attachment; filename="{download_name}"'
    return response


# ... all your imports ...

def save_chat_message(request, sender, content):
    """Save a chat message (user or assistant) to session."""
    chat_history = request.session.get('chat_history', [])
    chat_history.append({
        'type': 'chat_message',
        'sender': sender,
        'content': content,
        'timestamp': datetime.utcnow().isoformat()
    })
    request.session['chat_history'] = chat_history
    request.session.modified = True  # Ensure session is saved

def save_visualization(request, chart_title, image_url, explanation):
    """Save a visualization block to session."""
    visualizations = request.session.get('visualizations', [])
    visualizations.append({
        'type': 'visualization',
        'chart_title': chart_title,
        'image_url': image_url,
        'explanation': explanation,
        'timestamp': datetime.utcnow().isoformat()
    })
    request.session['visualizations'] = visualizations
    request.session.modified = True  # Ensure session is saved

def save_query_result(request, user_input, sql_query, result_data, visualization_suggestion):
    """Save a query history block to session."""
    chat_history = request.session.get('chat_history', [])
    chat_history.append({
        'type': 'query',
        'user_input': user_input,
        'sql_query': sql_query,
        'result_data': result_data,  # Save full result data for table generation
        'result_summary': {
            'rows': len(result_data),
            'columns': list(result_data[0].keys()) if result_data else [],
        },
        'visualization_suggestion': visualization_suggestion,
        'timestamp': datetime.utcnow().isoformat()
    })
    request.session['chat_history'] = chat_history
    request.session.modified = True  # Ensure session is saved

def save_dataset_schema(request, dataset_name, data_profile):
    """Save dataset schema information to session for reports."""
    if 'dataset_schemas' not in request.session:
        request.session['dataset_schemas'] = []
    
    schema_info = {
        'type': 'dataset_schema',
        'dataset_name': dataset_name,
        'shape': data_profile['shape'],
        'columns': data_profile['columns'][:10],  # Limit to first 10 columns for space
        'data_health': data_profile['data_health'],
        'total_columns': len(data_profile['columns']),
        'timestamp': datetime.utcnow().isoformat()
    }
    
    # Only keep the latest schema (replace if exists)
    schemas = request.session.get('dataset_schemas', [])
    # Remove any existing schema for this dataset
    schemas = [s for s in schemas if s.get('dataset_name') != dataset_name]
    schemas.append(schema_info)
    
    request.session['dataset_schemas'] = schemas
    request.session.modified = True

def save_query_table(request, query_description, result_data, user_input="", sql_query=""):
    """Save query result table for report inclusion."""
    if 'query_tables' not in request.session:
        request.session['query_tables'] = []
    
    table_info = {
        'type': 'query_table',
        'title': query_description,
        'user_input': user_input,
        'sql_query': sql_query,
        'data': result_data,
        'row_count': len(result_data),
        'columns': list(result_data[0].keys()) if result_data else [],
        'timestamp': datetime.utcnow().isoformat()
    }
    
    tables = request.session.get('query_tables', [])
    tables.append(table_info)
    request.session['query_tables'] = tables
    request.session.modified = True



def report_builder(request):
    """Enhanced report builder with complete chat history and editing capabilities"""
    # Start with an empty canvas: do not auto-populate from session
    blocks = []
    
    # Get dataset info for context
    dataset_display_name = request.session.get("dataset_display_name", "Dataset")
    
    return render(request, "enhanced_report_builder.html", {
        "blocks": blocks,
        "dataset_name": dataset_display_name,
        "total_blocks": len(blocks),
    })


def suggest_visualization_type(df, query):
    """Suggest the best visualization type based on data characteristics"""
    query_lower = query.lower()
    
    # Check for aggregation functions
    if any(func in query_lower for func in ['count', 'sum', 'avg', 'average', 'max', 'min']):
        if 'group by' in query_lower:
            return "bar_chart"
        elif len(df) == 1:
            return "metric_card"
        else:
            return "bar_chart"
    
    # Check for time series
    date_columns = [col for col in df.columns if any(word in col.lower() for word in ['date', 'time', 'year', 'month'])]
    if date_columns and len(df) > 1:
        return "line_chart"


def chart_builder(request):
    """New Chart.js based chart builder with axis selection"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            action = data.get('action')
            
            print(f"Chart builder action: {action}")  # Debug log
            
            if action == 'get_columns':
                # Return available columns for axis selection
                query_results = data.get('query_results', [])
                columns = data.get('columns', [])
                
                print(f"Query results length: {len(query_results)}")  # Debug log
                print(f"Columns: {columns}")  # Debug log
                
                if not query_results or not columns:
                    return JsonResponse({'error': 'No data provided'}, status=400)
                
                df = pd.DataFrame(query_results, columns=columns)
                
                # Categorize columns
                numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
                categorical_columns = df.select_dtypes(include=['object']).columns.tolist()
                datetime_columns = []
                
                # Try to detect datetime columns
                for col in categorical_columns[:]:  # Use slice to avoid modifying list while iterating
                    try:
                        sample_data = df[col].dropna().iloc[:5]
                        if len(sample_data) > 0:
                            pd.to_datetime(sample_data)
                            datetime_columns.append(col)
                            categorical_columns.remove(col)
                    except:
                        pass
                
                column_info = {
                    'numeric': numeric_columns,
                    'categorical': categorical_columns,
                    'datetime': datetime_columns,
                    'all': columns
                }
                
                print(f"Column info: {column_info}")  # Debug log
                
                return JsonResponse({
                    'success': True,
                    'columns': column_info
                })
            
            elif action == 'generate_chart':
                # Generate Chart.js configuration
                chart_type = data.get('chart_type', 'bar')
                x_axis = data.get('x_axis', [])
                y_axis = data.get('y_axis', [])
                query_results = data.get('query_results', [])
                columns = data.get('columns', [])
                
                if not query_results or not columns:
                    return JsonResponse({'error': 'No data provided'}, status=400)
                
                if not x_axis or not y_axis:
                    return JsonResponse({'error': 'Please select both X and Y axis variables'}, status=400)
                
                df = pd.DataFrame(query_results, columns=columns)
                
                # Generate Chart.js configuration
                chart_config = generate_chartjs_config(df, chart_type, x_axis, y_axis)
                
                return JsonResponse({
                    'success': True,
                    'chart_config': chart_config
                })
                
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)
    
    return JsonResponse({'error': 'Method not allowed'}, status=405)


def generate_chartjs_config(df, chart_type, x_axis, y_axis):
    """Generate Chart.js configuration based on selected axes and chart type"""
    
    # Handle multiple X and Y axes
    x_col = x_axis[0] if isinstance(x_axis, list) else x_axis
    y_cols = y_axis if isinstance(y_axis, list) else [y_axis]
    
    # Prepare data based on chart type
    if chart_type in ['bar', 'line']:
        # Group data if needed
        if df[x_col].dtype == 'object':  # Categorical X-axis
            grouped = df.groupby(x_col)[y_cols].agg('mean').reset_index()
            labels = grouped[x_col].tolist()
            
            datasets = []
            colors = ['#FF6384', '#36A2EB', '#FFCE56', '#4BC0C0', '#9966FF', '#FF9F40', '#FF6384', '#C9CBCF']
            
            for i, y_col in enumerate(y_cols):
                datasets.append({
                    'label': y_col,
                    'data': grouped[y_col].tolist(),
                    'backgroundColor': colors[i % len(colors)] if chart_type == 'bar' else 'transparent',
                    'borderColor': colors[i % len(colors)],
                    'borderWidth': 2,
                    'fill': False
                })
        else:  # Numeric X-axis
            labels = df[x_col].tolist()
            datasets = []
            colors = ['#FF6384', '#36A2EB', '#FFCE56', '#4BC0C0', '#9966FF', '#FF9F40']
            
            for i, y_col in enumerate(y_cols):
                datasets.append({
                    'label': y_col,
                    'data': df[y_col].tolist(),
                    'backgroundColor': colors[i % len(colors)] if chart_type == 'bar' else 'transparent',
                    'borderColor': colors[i % len(colors)],
                    'borderWidth': 2,
                    'fill': False
                })
    
    elif chart_type == 'scatter':
        datasets = []
        colors = ['#FF6384', '#36A2EB', '#FFCE56', '#4BC0C0', '#9966FF', '#FF9F40']
        
        for i, y_col in enumerate(y_cols):
            scatter_data = [{'x': x, 'y': y} for x, y in zip(df[x_col], df[y_col])]
            datasets.append({
                'label': f'{y_col} vs {x_col}',
                'data': scatter_data,
                'backgroundColor': colors[i % len(colors)],
                'borderColor': colors[i % len(colors)],
            })
        labels = []
    
    elif chart_type == 'pie':
        # For pie charts, use first Y column and group by X
        if df[x_col].dtype == 'object':
            grouped = df.groupby(x_col)[y_cols[0]].sum().reset_index()
            labels = grouped[x_col].tolist()
            data = grouped[y_cols[0]].tolist()
        else:
            # If X is numeric, create bins
            df[f'{x_col}_binned'] = pd.cut(df[x_col], bins=5, precision=0)
            grouped = df.groupby(f'{x_col}_binned')[y_cols[0]].sum().reset_index()
            labels = [str(interval) for interval in grouped[f'{x_col}_binned']]
            data = grouped[y_cols[0]].tolist()
        
        datasets = [{
            'data': data,
            'backgroundColor': ['#FF6384', '#36A2EB', '#FFCE56', '#4BC0C0', '#9966FF', '#FF9F40', '#FF6384', '#C9CBCF']
        }]
    
    # Chart.js configuration
    config = {
        'type': chart_type,
        'data': {
            'labels': labels,
            'datasets': datasets
        },
        'options': {
            'responsive': True,
            'maintainAspectRatio': False,
            'scales': {},
            'plugins': {
                'legend': {
                    'display': True,
                    'position': 'top'
                },
                'title': {
                    'display': True,
                    'text': f'{chart_type.title()} Chart: {", ".join(y_cols)} vs {x_col}'
                }
            }
        }
    }
    
    # Add scales for non-pie charts
    if chart_type != 'pie':
        config['options']['scales'] = {
            'x': {
                'display': True,
                'title': {
                    'display': True,
                    'text': x_col
                }
            },
            'y': {
                'display': True,
                'title': {
                    'display': True,
                    'text': ', '.join(y_cols)
                }
            }
        }
    
    return config


def generate_chart(request):
    """Legacy chart generation - redirects to new chart builder"""
    return JsonResponse({'error': 'Please use the new chart builder interface'}, status=400)


def generate_plotly_config(df, chart_type, x_axis, y_axis):
    """
    Generate Plotly.js configuration as a JSON object based on selected axes and chart type.
    """
    # Handle multiple X and Y axes
    x_col = x_axis[0] if isinstance(x_axis, list) else x_axis
    y_cols = y_axis if isinstance(y_axis, list) else [y_axis]

    fig = go.Figure()

    if chart_type == 'bar':
        # Group data for bar chart if x-axis is categorical
        if df[x_col].dtype == 'object':
            grouped_df = df.groupby(x_col)[y_cols].mean().reset_index()
            for y_col in y_cols:
                fig.add_trace(go.Bar(x=grouped_df[x_col], y=grouped_df[y_col], name=y_col))
        else:
            for y_col in y_cols:
                fig.add_trace(go.Bar(x=df[x_col], y=df[y_col], name=y_col))
        
        fig.update_layout(title=f'Bar Chart: {", ".join(y_cols)} vs {x_col}', xaxis_title=x_col, yaxis_title=", ".join(y_cols))
        
    elif chart_type == 'line':
        for y_col in y_cols:
            fig.add_trace(go.Scatter(x=df[x_col], y=df[y_col], mode='lines+markers', name=y_col))
        
        fig.update_layout(title=f'Line Chart: {", ".join(y_cols)} vs {x_col}', xaxis_title=x_col, yaxis_title=", ".join(y_cols))
        
    elif chart_type == 'scatter':
        for y_col in y_cols:
            fig.add_trace(go.Scatter(x=df[x_col], y=df[y_col], mode='markers', name=y_col))
        
        fig.update_layout(title=f'Scatter Plot: {", ".join(y_cols)} vs {x_col}', xaxis_title=x_col, yaxis_title=", ".join(y_cols))
        
    elif chart_type == 'pie':
        # Plotly.py supports pie charts natively
        if df[x_col].dtype == 'object':
            grouped = df.groupby(x_col)[y_cols[0]].sum().reset_index()
            fig = px.pie(grouped, values=y_cols[0], names=x_col, title=f'Pie Chart: {y_cols[0]} by {x_col}')
        else:
            # For numeric x-axis, create bins and then plot
            df[f'{x_col}_binned'] = pd.cut(df[x_col], bins=5, precision=0)
            grouped = df.groupby(f'{x_col}_binned')[y_cols[0]].sum().reset_index()
            grouped[f'{x_col}_binned'] = grouped[f'{x_col}_binned'].astype(str)
            fig = px.pie(grouped, values=y_cols[0], names=f'{x_col}_binned', title=f'Pie Chart: {y_cols[0]} by {x_col}')

    elif chart_type == 'heatmap':
        # Heatmaps require a Z-axis, usually derived from a pivot table
        # We'll assume y_cols has the row and column names, and a value column
        if len(y_cols) < 2:
            return {'error': 'Heatmap requires at least two Y-axis columns and a value column.'}
        
        # We need three columns: x, y, and z (value)
        z_col = y_cols[0]
        y_col_heatmap = y_cols[1]
        
        fig = go.Figure(data=go.Heatmap(
                z=df[z_col],
                x=df[x_col],
                y=df[y_col_heatmap],
                colorscale='Viridis'))
        
        fig.update_layout(
            title=f'Heatmap: {z_col} vs {x_col} and {y_col_heatmap}',
            xaxis_title=x_col,
            yaxis_title=y_col_heatmap)
            
    elif chart_type == 'bubble':
        # Bubble charts require a size parameter, which we'll assume is the second y-column
        if len(y_cols) < 2:
            return {'error': 'Bubble chart requires at least two Y-axis columns for value and size.'}
        
        fig = go.Figure(data=[go.Scatter(
            x=df[x_col],
            y=df[y_cols[0]],
            mode='markers',
            marker=dict(
                size=df[y_cols[1]], # Use the second y-column for bubble size
                sizemode='area',
                sizeref=2.*max(df[y_cols[1]])/(40.**2), # Adjust sizing
                sizemin=4
            )
        )])
        
        fig.update_layout(title=f'Bubble Chart: {y_cols[0]} vs {x_col} (Size by {y_cols[1]})',
                          xaxis_title=x_col,
                          yaxis_title=y_cols[0])

    else:
        return {'error': 'Unsupported chart type'}
    
    # Convert the figure to a JSON serializable object
    return json.loads(fig.to_json())

def session_cleanup_endpoint(request):
    """Endpoint to handle complete session cleanup including Supabase data"""
    if request.method == 'POST':
        try:
            # Perform complete cleanup
            success = complete_session_cleanup(request.session)
            
            if success:
                return JsonResponse({
                    'status': 'success',
                    'message': 'Session cleaned up successfully, all data removed from Supabase'
                })
            else:
                return JsonResponse({
                    'status': 'error',
                    'message': 'Failed to complete session cleanup'
                })
                
        except Exception as e:
            return JsonResponse({
                'status': 'error',
                'message': f'Error during cleanup: {str(e)}'
            })
    
    return JsonResponse({'status': 'error', 'message': 'Invalid request method'})


def remove_dataset(request, dataset_index):
    """Remove a dataset from the current session (and optionally its uploaded file)."""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'})

    try:
        datasets_info = request.session.get('datasets', [])
        if not (0 <= dataset_index < len(datasets_info)):
            return JsonResponse({'status': 'error', 'message': 'Invalid dataset index'})

        dataset_info = datasets_info[dataset_index]
        path_in_bucket = dataset_info.get('supabase_path')

        # Remove dataset entry from session list
        del datasets_info[dataset_index]
        request.session['datasets'] = datasets_info

        # Adjust active dataset index
        active_index = request.session.get('active_dataset_index', 0)
        if len(datasets_info) == 0:
            request.session['active_dataset_index'] = 0
            request.session['current_dataset_path'] = None
        else:
            if dataset_index < active_index:
                active_index -= 1
            elif dataset_index == active_index:
                # Move to nearest valid index
                active_index = min(active_index, len(datasets_info) - 1)
            request.session['active_dataset_index'] = active_index
            # Update backward-compat path
            request.session['current_dataset_path'] = datasets_info[active_index]['supabase_path']

        request.session.modified = True

        # Try to delete the file from Supabase storage (best-effort)
        try:
            if path_in_bucket:
                supabase.storage.from_(SUPABASE_BUCKET).remove([path_in_bucket])
                print(f"🗑️ Deleted dataset file from Supabase: {path_in_bucket}")
        except Exception as del_err:
            print(f"⚠️ Could not delete dataset file: {del_err}")

        return JsonResponse({
            'status': 'success',
            'remaining': len(datasets_info),
            'active_index': request.session.get('active_dataset_index', 0)
        })

    except Exception as e:
        print(f"❌ remove_dataset failed: {e}")
        return JsonResponse({'status': 'error', 'message': str(e)})

def session_end_cleanup(request):
    """Handle cleanup when session ends (browser close, timeout, etc.)"""
    try:
        # This will be called by JavaScript when the page is about to unload
        complete_session_cleanup(request.session)
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)})

def advanced_visualization_builder(request, dataset_name=None):
    """Advanced visualization builder with adaptive interface"""
    try:
        # Debug session info
        print(f"DEBUG: Session keys: {list(request.session.keys())}")
        print(f"DEBUG: Session datasets: {request.session.get('datasets', [])}")
        print(f"DEBUG: Active dataset index: {request.session.get('active_dataset_index', 0)}")
        print(f"DEBUG: Session ID: {request.session.session_key}")
        
        # Get dataset information from active dataset
        dataset_path, dataset_display_name = get_active_dataset_path(request)
        if not dataset_path:
            print("DEBUG: No active dataset found in session")
            messages.error(request, "No dataset found. Please upload a dataset first.")
            return redirect("index")
        
        # Load dataset from Supabase
        try:
            print(f"DEBUG: Loading dataset from: {dataset_path}")
            res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_path)
            df = pd.read_csv(io.BytesIO(res))
            print(f"DEBUG: Dataset loaded successfully, shape: {df.shape}")
        except Exception as e:
            print(f"DEBUG: Error loading dataset: {e}")
            
            # Try to list what files are actually available
            try:
                files = supabase.storage.from_(SUPABASE_BUCKET).list("sessions/")
                print(f"DEBUG: Available files: {[f.get('name', 'unknown') for f in files[:10]]}")
            except Exception as list_error:
                print(f"DEBUG: Cannot list files: {list_error}")
            
            messages.error(request, "Failed to load dataset. The file may have expired. Please upload the dataset again.")
            return redirect("index")
        
        # Get column information
        numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
        categorical_columns = df.select_dtypes(include=['object', 'category']).columns.tolist()
        datetime_columns = df.select_dtypes(include=['datetime64']).columns.tolist()
        all_columns = df.columns.tolist()
        
        # Get sample data for preview
        sample_data = df.head(10).to_dict('records')
        
        context = {
            'dataset_name': dataset_display_name,
            'dataset_internal_name': dataset_path.split('/')[-1] if dataset_path else "",
            'columns': all_columns,
            'numeric_columns': numeric_columns,
            'categorical_columns': categorical_columns,
            'datetime_columns': datetime_columns,
            'sample_data': sample_data,
            'row_count': len(df),
            'column_count': len(df.columns)
        }
        
        return render(request, 'advanced_visualization_builder_new.html', context)
        
    except Exception as e:
        print(f"Error in advanced_visualization_builder: {e}")
        messages.error(request, f"Error loading visualization builder: {str(e)}")
        return redirect("index")

def generate_advanced_chart(request):
    """Generate charts using Matplotlib and Seaborn based on user selections"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Invalid request method'})
    
    try:
        import matplotlib
        matplotlib.use('Agg')  # Use non-interactive backend
        import matplotlib.pyplot as plt
        import seaborn as sns
        from matplotlib.colors import ListedColormap
        import warnings
        warnings.filterwarnings('ignore')
        
        # Get request data
        data = json.loads(request.body)
        chart_type = data.get('chart_type')
        chart_config = data.get('config', {})
        
        # Debug logging
        print(f"Chart type: {chart_type}")
        print(f"Chart config: {chart_config}")
        
        # Load dataset from active dataset
        dataset_path, display_name = get_active_dataset_path(request)
        print(f"DEBUG: Session keys: {list(request.session.keys())}")
        print(f"DEBUG: Active dataset: {display_name}")
        print(f"DEBUG: Active dataset path: {dataset_path}")
        print(f"DEBUG: Session ID: {request.session.session_key}")
        
        if not dataset_path:
            print("DEBUG: No active dataset found in session")
            return JsonResponse({'error': 'No dataset found in session. Please upload a dataset first.'})
        
        try:
            print(f"DEBUG: Attempting to download from Supabase: {dataset_path}")
            print(f"DEBUG: Bucket: {SUPABASE_BUCKET}")
            res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_path)
            df = pd.read_csv(io.BytesIO(res))
            print(f"DEBUG: Successfully loaded dataset with shape: {df.shape}")
            print(f"DEBUG: Columns: {df.columns.tolist()}")
        except Exception as e:
            print(f"DEBUG: Error loading dataset from Supabase: {e}")
            print(f"DEBUG: Dataset path: {dataset_path}")
            print(f"DEBUG: Error type: {type(e)}")
            
            # Try to list files in the bucket to see what's available
            try:
                files = supabase.storage.from_(SUPABASE_BUCKET).list("sessions/")
                print(f"DEBUG: Available files in bucket: {[f.get('name') for f in files[:5]]}")
            except Exception as list_error:
                print(f"DEBUG: Could not list files: {list_error}")
                
            return JsonResponse({'error': f'Failed to load dataset: {str(e)}. Please try uploading the dataset again.'})
        
        # Validate dataset
        if df.empty:
            return JsonResponse({'error': 'Dataset is empty'})
        
        # Set style and color palette
        plt.style.use('dark_background')
        selected_palette = chart_config.get('palette', 'husl')
        sns.set_palette(selected_palette)
        
        # Process configuration for specific chart types
        if chart_type == 'pie':
            # Determine subtype if not explicitly provided
            subtype = chart_config.get('chart_subtype')
            if not subtype:
                if chart_config.get('group_by_column') and chart_config.get('category_column'):
                    subtype = 'grouped'
                elif chart_config.get('values_column'):
                    subtype = 'values'
                else:
                    subtype = 'simple'
            chart_config['chart_subtype'] = subtype

            if subtype == 'simple':
                chart_config['simple_column'] = chart_config.get('category_column')
                # Optional value column for sum by category
                if chart_config.get('value_column'):
                    chart_config['value_column'] = chart_config.get('value_column')
            elif subtype == 'grouped':
                # Expect 'group_by_column' and 'category_column'
                chart_config['group_by_column'] = chart_config.get('group_by_column')
                chart_config['category_column'] = chart_config.get('category_column')
                chart_config['aggregation_method'] = chart_config.get('aggregation_method', 'count')
                chart_config['specific_group'] = chart_config.get('specific_group')
            elif subtype == 'values':
                chart_config['values_column'] = chart_config.get('values_column')
                chart_config['labels_column'] = chart_config.get('labels_column')
        
        # Create figure
        plt.figure(figsize=(12, 8))
        
        # Generate chart based on type
        if chart_type == 'bar':
            chart_html = generate_bar_chart(df, chart_config)
        elif chart_type == 'line':
            chart_html = generate_line_chart(df, chart_config)
        elif chart_type == 'scatter':
            chart_html = generate_scatter_chart(df, chart_config)
        elif chart_type == 'histogram':
            chart_html = generate_histogram_chart(df, chart_config)
        elif chart_type == 'box':
            chart_html = generate_box_chart(df, chart_config)
        elif chart_type == 'pie':
            chart_html = generate_pie_chart(df, chart_config)
        elif chart_type == 'heatmap':
            chart_html = generate_heatmap_chart(df, chart_config)
        elif chart_type == 'violin':
            chart_html = generate_violin_chart(df, chart_config)
        elif chart_type == 'area':
            chart_html = generate_area_chart(df, chart_config)
        elif chart_type == 'density':
            chart_html = generate_density_chart(df, chart_config)
        else:
            return JsonResponse({'error': f'Unsupported chart type: {chart_type}'})
        
        # Clear the plot
        plt.close('all')
        
        # Store in session for report builder
        if 'visualizations' not in request.session:
            request.session['visualizations'] = []
        
        # Extract base64 image from HTML for report compatibility
        image_url = ""
        if 'data:image/png;base64,' in chart_html:
            # Extract the data URL from the img tag
            import re
            match = re.search(r'src="(data:image/png;base64,[^"]+)"', chart_html)
            if match:
                image_url = match.group(1)
        
        chart_info = {
            'type': 'visualization',  # Match the type used by auto-generated charts
            'chart_title': f"{chart_type.title()} Chart: {chart_config.get('title', '')}",
            'image_url': image_url,  # Use image_url for PDF compatibility
            'html': chart_html,  # Keep HTML for web display
            'explanation': f"Advanced {chart_type} chart generated with custom configuration",
            'chart_type': chart_type,
            'config': chart_config,
            'timestamp': datetime.now().isoformat()
        }
        
        request.session['visualizations'].append(chart_info)
        request.session.modified = True
        
        return JsonResponse({
            'success': True,
            'chart_html': chart_html,
            'chart_info': chart_info
        })
        
    except Exception as e:
        plt.close('all')  # Cleanup on error
        return JsonResponse({'error': f'Error generating chart: {str(e)}'})

def generate_bar_chart(df, config):
    """Generate bar chart with Matplotlib/Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        x_col = config.get('x_column')
        y_cols = config.get('y_columns', [])
        title = config.get('title', 'Bar Chart')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        stacked = config.get('stacked', False)
        horizontal = config.get('horizontal', False)
        alpha = float(config.get('alpha', 0.8))
        edge_color = config.get('edge_color') or None
        line_width = float(config.get('line_width', 0))
        bar_width = float(config.get('bar_width', 0.8))
        
        if not x_col or not y_cols:
            raise ValueError("X column and Y columns are required")
        
        # Handle categorical data
        if df[x_col].dtype == 'object':
            # Group by categorical column and aggregate
            if len(y_cols) == 1:
                grouped_df = df.groupby(x_col)[y_cols[0]].sum().reset_index()
                x_data = grouped_df[x_col]
                y_data = grouped_df[y_cols[0]]
                
                color = sns.color_palette(color_palette, 1)[0]
                if horizontal:
                    plt.barh(x_data, y_data, color=color, alpha=alpha, edgecolor=edge_color, linewidth=line_width, height=bar_width)
                    plt.xlabel(y_cols[0])
                    plt.ylabel(x_col)
                else:
                    plt.bar(x_data, y_data, color=color, alpha=alpha, edgecolor=edge_color, linewidth=line_width, width=bar_width)
                    plt.xlabel(x_col)
                    plt.ylabel(y_cols[0])
            else:
                # Multiple y columns
                grouped_df = df.groupby(x_col)[y_cols].sum()
                
                colors = sns.color_palette(color_palette, len(y_cols))
                ax = plt.gca()
                grouped_df.plot(kind='barh' if horizontal else 'bar',
                                 stacked=stacked,
                                 color=colors,
                                 ax=ax,
                                 alpha=alpha)
                # Apply edgecolor/linewidth to bars
                for patch in ax.patches:
                    if edge_color is not None:
                        patch.set_edgecolor(edge_color)
                    patch.set_linewidth(line_width)
        else:
            # Numeric x column
            for i, y_col in enumerate(y_cols):
                color = sns.color_palette(color_palette, len(y_cols))[i]
                if horizontal:
                    plt.barh(df[x_col], df[y_col], alpha=alpha, label=y_col, color=color, edgecolor=edge_color, linewidth=line_width, height=bar_width)
                else:
                    plt.bar(df[x_col], df[y_col], alpha=alpha, label=y_col, color=color, edgecolor=edge_color, linewidth=line_width, width=bar_width)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xticks(rotation=45)
        # Legend positioning
        legend_pos = config.get('legend_position', 'best')
        if legend_pos and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating bar chart: {str(e)}")

def generate_line_chart(df, config):
    """Generate line chart with Matplotlib/Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        x_col = config.get('x_column')
        y_cols = config.get('y_columns', [])
        title = config.get('title', 'Line Chart')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        markers = config.get('markers', True)
        marker = config.get('marker', 'o')
        line_style = config.get('line_style', '-')
        line_width = float(config.get('line_width', 2))
        alpha = float(config.get('alpha', 1))
        
        if not x_col or not y_cols:
            raise ValueError("X column and Y columns are required")
        
        colors = sns.color_palette(color_palette, len(y_cols))
        
        for i, y_col in enumerate(y_cols):
            if markers:
                plt.plot(df[x_col], df[y_col], marker=marker,
                        linestyle=line_style, color=colors[i], 
                        label=y_col, linewidth=line_width, markersize=4, alpha=alpha)
            else:
                plt.plot(df[x_col], df[y_col], 
                        linestyle=line_style, color=colors[i], 
                        label=y_col, linewidth=line_width, alpha=alpha)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_col)
        plt.ylabel(' / '.join(y_cols))
        # Legend positioning
        legend_pos = config.get('legend_position', 'best')
        if legend_pos and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.xticks(rotation=45)
        plt.grid(True, alpha=0.3)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating line chart: {str(e)}")

def generate_scatter_chart(df, config):
    """Generate scatter plot with Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        x_col = config.get('x_column')
        y_col = config.get('y_column')
        hue_col = config.get('hue_column')
        size_col = config.get('size_column')
        style_col = config.get('style_column')
        title = config.get('title', 'Scatter Plot')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        marker = config.get('marker', 'o')
        alpha = float(config.get('alpha', 0.7))
        edge_color = config.get('edge_color') or '#000000'
        line_width = float(config.get('line_width', 0))
        point_size = float(config.get('point_size', 60))
        color_array_col = config.get('color_array_column')
        colormap = config.get('colormap', 'viridis')
        
        if not x_col or not y_col:
            raise ValueError("X and Y columns are required")
        
        # Create scatter plot
        if color_array_col:
            # Use Matplotlib to map numeric values to colors via colormap
            x_vals = df[x_col]
            y_vals = df[y_col]
            c_vals = df[color_array_col]
            # Handle sizes: column-based or constant
            if size_col:
                # Normalize size column to a reasonable range
                sizes = df[size_col].astype(float)
                # Avoid negative/zero sizes
                sizes = sizes.clip(lower=1)
                # Scale sizes to around the requested point_size
                scale = point_size / (sizes.mean() if sizes.mean() != 0 else 1.0)
                s_param = sizes * scale
            else:
                s_param = point_size

            plt.scatter(x_vals, y_vals, c=c_vals, cmap=colormap,
                        s=s_param, marker=marker, alpha=alpha,
                        edgecolors=edge_color, linewidths=line_width)
        else:
            # Use Seaborn with hue/style/size mapping and palette
            scatter_kwargs = {
                'data': df,
                'x': x_col,
                'y': y_col,
                'hue': hue_col,
                'style': style_col,
                'palette': color_palette,
                'alpha': alpha,
                'marker': marker,
                'edgecolor': edge_color,
                'linewidth': line_width
            }
            if size_col:
                scatter_kwargs['size'] = size_col
            else:
                scatter_kwargs['s'] = point_size

            sns.scatterplot(**scatter_kwargs)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_col)
        plt.ylabel(y_col)
        # Legend repositioning/hiding for seaborn legends
        legend_pos = config.get('legend_position', 'best')
        ax = plt.gca()
        leg = ax.get_legend()
        if legend_pos == 'none':
            if leg:
                leg.remove()
        elif leg:
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating scatter chart: {str(e)}")

def generate_histogram_chart(df, config):
    """Generate histogram with Matplotlib/Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        columns = config.get('columns', [])
        title = config.get('title', 'Histogram')
        bins = config.get('bins', 30)
        density = config.get('density', False)
        alpha = float(config.get('alpha', 0.7))
        edge_color = config.get('edge_color') or None
        line_width = float(config.get('line_width', 0))
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        
        if not columns:
            raise ValueError("At least one column is required")
        
        colors = sns.color_palette(color_palette, len(columns))
        
        for i, col in enumerate(columns):
                plt.hist(df[col].dropna(), bins=bins, alpha=alpha, 
                    label=col, color=colors[i], density=density,
                    edgecolor=edge_color, linewidth=line_width)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel('Value')
        plt.ylabel('Density' if density else 'Frequency')
        # Legend positioning
        legend_pos = config.get('legend_position', 'best')
        if legend_pos and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.grid(True, alpha=0.3)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating histogram: {str(e)}")

def generate_box_chart(df, config):
    """Generate box plot with Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        y_cols = config.get('y_columns', [])
        x_col = config.get('x_column')  # Optional grouping column
        title = config.get('title', 'Box Plot')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        notch = config.get('notch', False)
        show_fliers = config.get('show_fliers', True)
        line_width = float(config.get('line_width', 1))
        
        if not y_cols:
            raise ValueError("At least one Y column is required")
        
        if x_col:
            # Grouped box plot
            for y_col in y_cols:
                sns.boxplot(data=df, x=x_col, y=y_col, palette=color_palette, showfliers=show_fliers)
                for artist in plt.gca().artists:
                    artist.set_linewidth(line_width)
                plt.xticks(rotation=45)
        else:
            # Simple box plot
            data_to_plot = [df[col].dropna() for col in y_cols]
            plt.boxplot(
                data_to_plot,
                labels=y_cols,
                patch_artist=True,
                notch=notch,
                showfliers=show_fliers,
                boxprops=dict(facecolor=sns.color_palette(color_palette, 1)[0], linewidth=line_width),
                whiskerprops=dict(linewidth=line_width),
                capprops=dict(linewidth=line_width),
                medianprops=dict(linewidth=line_width)
            )
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating box plot: {str(e)}")

def generate_pie_chart(df, config):
    """Generate advanced pie chart with complex grouping and filtering capabilities"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        chart_subtype = config.get('chart_subtype', 'simple')
        title = config.get('title', 'Pie Chart')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        max_slices = config.get('max_slices', 8)
        donut = config.get('donut', False)
        explode_factor = float(config.get('explode', 0))
        
        # Debug logging
        print(f"Pie chart subtype: {chart_subtype}")
        print(f"Config: {config}")
        print(f"Available columns: {df.columns.tolist()}")
        
        # Create figure with proper size
        plt.figure(figsize=(12, 8))
        
        if chart_subtype == 'simple':
            # Simple distribution - count occurrences in a single column
            simple_column = config.get('simple_column') or config.get('category_column')
            if not simple_column or simple_column not in df.columns:
                raise ValueError(f"Category column '{simple_column}' is required and must exist in dataset")
            
            # Check if there's a value column to sum by
            value_column = config.get('value_column')
            if value_column and value_column in df.columns:
                # Sum values by category
                grouped_data = df.groupby(simple_column)[value_column].sum().sort_values(ascending=False)
                value_counts = grouped_data.head(max_slices)
                labels = value_counts.index.tolist()
                values = value_counts.values
                
                if not title or title == 'Pie Chart':
                    title = f"Sum of {value_column} by {simple_column}"
            else:
                # Count occurrences by category
                value_counts = df[simple_column].value_counts().head(max_slices)
                labels = value_counts.index.tolist()
                values = value_counts.values
                
                if not title or title == 'Pie Chart':
                    title = f"Distribution of {simple_column}"
                
        elif chart_subtype == 'grouped':
            # Grouped distribution - e.g., Gender distribution by Department
            group_by_column = config.get('group_by_column')
            category_column = config.get('category_column')
            specific_group = config.get('specific_group')
            aggregation_method = config.get('aggregation_method', 'count')
            
            if not group_by_column or group_by_column not in df.columns:
                raise ValueError("Group by column is required and must exist in dataset")
            if not category_column or category_column not in df.columns:
                raise ValueError("Category column is required and must exist in dataset")
            
            # Filter to specific group if requested
            if specific_group:
                filtered_df = df[df[group_by_column] == specific_group]
                if len(filtered_df) == 0:
                    raise ValueError(f"No data found for group '{specific_group}'")
                
                # Count categories within the specific group
                value_counts = filtered_df[category_column].value_counts()
                labels = value_counts.index.tolist()
                values = value_counts.values
                
                if not title:
                    title = f"{category_column} Distribution in {specific_group}"
            else:
                # Show distribution across all groups
                if aggregation_method == 'count':
                    # Total count across all groups
                    value_counts = df[category_column].value_counts().head(max_slices)
                    labels = value_counts.index.tolist()
                    values = value_counts.values
                    
                    if not title:
                        title = f"Overall {category_column} Distribution"
                else:
                    # Percentage within each group (complex calculation)
                    cross_tab = pd.crosstab(df[group_by_column], df[category_column])
                    percentage_df = cross_tab.div(cross_tab.sum(axis=1), axis=0) * 100
                    
                    # Flatten and get top combinations
                    flattened = []
                    for group in percentage_df.index:
                        for category in percentage_df.columns:
                            if not pd.isna(percentage_df.loc[group, category]):
                                flattened.append({
                                    'label': f"{category} ({group})",
                                    'value': percentage_df.loc[group, category]
                                })
                    
                    # Sort and take top entries
                    flattened.sort(key=lambda x: x['value'], reverse=True)
                    flattened = flattened[:max_slices]
                    
                    labels = [item['label'] for item in flattened]
                    values = np.array([item['value'] for item in flattened])
                    
                    if not title:
                        title = f"{category_column} Percentage by {group_by_column}"
                        
        elif chart_subtype == 'values':
            # Value-based pie chart - use numeric values
            values_column = config.get('values_column')
            labels_column = config.get('labels_column')
            
            if not values_column or values_column not in df.columns:
                raise ValueError("Values column is required and must exist in dataset")
            
            if labels_column and labels_column in df.columns:
                # Group by labels and sum values
                grouped_df = df.groupby(labels_column)[values_column].sum().sort_values(ascending=False)
                grouped_df = grouped_df.head(max_slices)
                labels = grouped_df.index.tolist()
                values = grouped_df.values
            else:
                # Use top values with generic labels
                top_values = df[values_column].dropna().nlargest(max_slices)
                values = top_values.values
                labels = [f"Item {i+1}" for i in range(len(values))]
                
            if not title:
                title = f"Distribution by {values_column}"
        else:
            raise ValueError(f"Unknown pie chart subtype: {chart_subtype}")
        
        # Ensure we have valid data
        if len(values) == 0:
            raise ValueError("No data available for pie chart")
        
        # Filter out zero or negative values for percentage-based charts
        if chart_subtype != 'grouped' or config.get('aggregation_method') != 'percentage':
            valid_mask = values > 0
            values = values[valid_mask]
            labels = [labels[i] for i in range(len(labels)) if valid_mask[i]]
        
        if len(values) == 0:
            raise ValueError("No positive values found for pie chart")
        
        # Add "Others" category if we have more data than max_slices
        if chart_subtype == 'simple':
            total_unique = df[config.get('simple_column')].nunique()
            if total_unique > max_slices:
                others_count = df[config.get('simple_column')].value_counts().iloc[max_slices:].sum()
                if others_count > 0:
                    labels.append('Others')
                    values = np.append(values, others_count)
        
        # Generate colors
        colors = sns.color_palette(color_palette, len(values))
        
        # Create pie chart
        explode = [explode_factor] * len(values) if explode_factor > 0 else None
        wedges, texts, autotexts = plt.pie(
            values,
            labels=labels,
            autopct='%1.1f%%',
            startangle=90,
            colors=colors,
            textprops={'fontsize': 10, 'color': 'white'},
            explode=explode,
            wedgeprops={'width': 0.4} if donut else None
        )
        
        # Improve text visibility
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(9)
        
        for text in texts:
            text.set_color('white')
            text.set_fontsize(10)
        
        plt.title(title, fontsize=16, fontweight='bold', color='white', pad=20)
        plt.axis('equal')
        
        # Legend handling based on config
        show_legend = bool(config.get('show_legend', True))
        legend_pos = config.get('legend_position', 'best')
        if show_legend and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(wedges, labels, title="Categories", loc='upper center', bbox_to_anchor=(0.5, 1.15), fontsize=10)
            else:
                plt.legend(wedges, labels, title="Categories", loc=legend_pos, fontsize=10)
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        plt.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        print(f"Error in generate_pie_chart: {str(e)}")
        raise Exception(f"Error generating pie chart: {str(e)}")

def generate_heatmap_chart(df, config):
    """Generate heatmap with Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import base64
        from io import BytesIO
        
        columns = config.get('columns', [])
        title = config.get('title', 'Heatmap')
        color_palette = config.get('palette') or config.get('color_palette', 'viridis')
        # Overrides from config
        cmap = config.get('colormap', color_palette)
        annot = bool(config.get('annot', True))
        fmt = config.get('fmt', '.2f')
        line_width = float(config.get('line_width', 0))
        line_color = config.get('line_color') or 'black'
        center = config.get('center', 0)
        
        if not columns:
            # Use all numeric columns
            numeric_df = df.select_dtypes(include=[np.number])
        else:
            numeric_df = df[columns]
        
        # Calculate correlation matrix
        corr_matrix = numeric_df.corr()
        
        # Create heatmap
        sns.heatmap(
            corr_matrix,
            annot=annot,
            cmap=cmap,
            center=center,
            square=True,
            fmt=fmt,
            linewidths=line_width,
            linecolor=line_color
        )
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating heatmap: {str(e)}")

def generate_violin_chart(df, config):
    """Generate violin plot with Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        y_cols = config.get('y_columns', [])
        x_col = config.get('x_column')
        title = config.get('title', 'Violin Plot')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        inner = config.get('inner', 'box')
        split = bool(config.get('split', False))
        scale = config.get('scale', 'area')
        bw = config.get('bw', None)
        line_width = float(config.get('line_width', 1))
        
        if not y_cols:
            raise ValueError("At least one Y column is required")
        
        if x_col:
            # Grouped violin plot
            for y_col in y_cols:
                sns.violinplot(
                    data=df,
                    x=x_col,
                    y=y_col,
                    palette=color_palette,
                    inner=inner if inner != 'none' else None,
                    scale=scale,
                    bw=bw,
                    linewidth=line_width
                )
                plt.xticks(rotation=45)
        else:
            # Simple violin plot
            data_to_plot = [df[col].dropna() for col in y_cols]
            sns.violinplot(
                data=data_to_plot,
                palette=color_palette,
                inner=inner if inner != 'none' else None,
                scale=scale,
                bw=bw,
                linewidth=line_width
            )
            plt.xticks(range(len(y_cols)), y_cols)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating violin plot: {str(e)}")

def generate_area_chart(df, config):
    """Generate area chart with Matplotlib"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        x_col = config.get('x_column')
        y_cols = config.get('y_columns', [])
        title = config.get('title', 'Area Chart')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        stacked = config.get('stacked', True)
        
        if not x_col or not y_cols:
            raise ValueError("X column and Y columns are required")
        
        colors = sns.color_palette(color_palette, len(y_cols))
        
        if stacked:
            plt.stackplot(df[x_col], *[df[col] for col in y_cols], 
                         labels=y_cols, colors=colors, alpha=0.7)
        else:
            for i, y_col in enumerate(y_cols):
                plt.fill_between(df[x_col], df[y_col], alpha=0.7, 
                               color=colors[i], label=y_col)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_col)
        plt.ylabel(' / '.join(y_cols))
        # Legend positioning
        legend_pos = config.get('legend_position', 'best')
        if legend_pos and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.xticks(rotation=45)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating area chart: {str(e)}")

def generate_density_chart(df, config):
    """Generate density plot with Seaborn"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        import numpy as np
        import pandas as pd
        import base64
        from io import BytesIO
        
        columns = config.get('columns', [])
        title = config.get('title', 'Density Plot')
        color_palette = config.get('palette') or config.get('color_palette', 'husl')
        
        if not columns:
            raise ValueError("At least one column is required")
        
        colors = sns.color_palette(color_palette, len(columns))
        
        for i, col in enumerate(columns):
            sns.kdeplot(data=df, x=col, color=colors[i], label=col, fill=True, alpha=0.6)
        
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel('Value')
        plt.ylabel('Density')
        # Legend positioning
        legend_pos = config.get('legend_position', 'best')
        if legend_pos and legend_pos != 'none':
            if legend_pos == 'top':
                plt.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15))
            else:
                plt.legend(loc=legend_pos)
        plt.tight_layout()
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight', 
                   facecolor='#1a1a1a', edgecolor='none', dpi=100)
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        buffer.close()
        
        return f'<img src="data:image/png;base64,{image_base64}" style="max-width: 100%; height: auto;">'
        
    except Exception as e:
        raise Exception(f"Error generating density plot: {str(e)}")


def get_column_values(request):
    """Get unique values for a specific column"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Invalid request method'})
    
    try:
        data = json.loads(request.body)
        column_name = data.get('column')
        
        if not column_name:
            return JsonResponse({'error': 'Column name is required'})
        
        # Load dataset from active dataset
        dataset_path, display_name = get_active_dataset_path(request)
        if not dataset_path:
            return JsonResponse({'error': 'No dataset found'})
        
        try:
            res = supabase.storage.from_(SUPABASE_BUCKET).download(dataset_path)
            df = pd.read_csv(io.BytesIO(res))
        except Exception as e:
            return JsonResponse({'error': 'Failed to load dataset'})
        
        if column_name not in df.columns:
            return JsonResponse({'error': f'Column {column_name} not found'})
        
        # Get unique values (limit to 50 for performance)
        unique_values = df[column_name].dropna().unique()[:50]
        unique_values = [str(val) for val in unique_values]
        
        return JsonResponse({'values': unique_values})
        
    except Exception as e:
        return JsonResponse({'error': str(e)})

def debug_session(request):
    """Debug session data"""
    context = {
        'session_keys': list(request.session.keys()),
        'dataset_path': request.session.get('dataset_path'),
        'dataset_display_name': request.session.get('dataset_display_name'),
    }
    return render(request, 'session_debug.html', context)


def transcribe_audio(request):
    """Transcribe audio using Groq Whisper API"""
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'message': 'Invalid request'}, status=400)

    try:
        # Get the audio file from the POST request
        audio_file = request.FILES.get('audio')
        if not audio_file:
            return JsonResponse({'status': 'error', 'message': 'No audio file received'}, status=400)

        print(f"🎤 Transcribing audio file: {audio_file.name}, size: {audio_file.size} bytes")

        # Use the Groq client to transcribe
        # We must send the file with a name, so we use a tuple
        transcription = client.audio.transcriptions.create(
            model="whisper-large-v3",
            file=("audio.webm", audio_file.read(), audio_file.content_type)
        )

        print(f"✅ Transcription successful: {transcription.text}")

        return JsonResponse({
            'status': 'success',
            'text': transcription.text
        })

    except Exception as e:
        print(f"❌ Transcription error: {e}")
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)